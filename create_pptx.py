#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""OTTOBITE Garson Davranışları ve İş Önceliği Rehberi - Modern PPTX Sunum Oluşturucu"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Renkler
KIREMIT = RGBColor(0xB8, 0x4C, 0x3C)
GRI_KUTU = RGBColor(0x3A, 0x3A, 0x50)  # Koyu gri-mor arka plan kutusu
BEYAZ = RGBColor(0xFF, 0xFF, 0xFF)
ACIK_GRI = RGBColor(0xCC, 0xCC, 0xCC)
ACCENT = RGBColor(0xC9, 0x5A, 0x48)
ARKA_PLAN = RGBColor(0x1A, 0x1A, 0x2E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_gradient_background(slide):
    """Koyu arka plan ekle"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = ARKA_PLAN
    background.line.fill.background()
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_title_shape(slide, title_text, is_main=False):
    """Kiremit renkli başlık ekle"""
    if is_main:
        left, top, width, height = Inches(0.5), Inches(0.4), Inches(12.333), Inches(1.2)
        font_size = Pt(36)
    else:
        left, top, width, height = Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
        font_size = Pt(26)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = font_size
    p.font.bold = True
    p.font.color.rgb = KIREMIT
    p.font.name = "Arial"
    return title_box

def add_content_box(slide, content_lines, start_top=1.3, base_font_size=16):
    """İçerik kutusu ekle - gri arka planlı kutular ve maddeler ile"""
    current_top = Inches(start_top)
    left = Inches(0.5)
    width = Inches(12.333)
    max_bottom = Inches(6.8)  # Sayfa altından en az 0.7 inch boşluk
    
    # İçerik miktarına göre font boyutu ayarla
    total_items = len([l for l in content_lines if l.strip()])
    if total_items > 15:
        font_size = base_font_size - 2
        line_height = 0.32
        box_height = 0.45
    elif total_items > 12:
        font_size = base_font_size - 1
        line_height = 0.35
        box_height = 0.48
    else:
        font_size = base_font_size
        line_height = 0.38
        box_height = 0.52
    
    for line in content_lines:
        line = line.strip()
        if not line:
            current_top += Inches(0.12)
            continue
        
        # Maksimum yükseklik kontrolü
        if current_top > max_bottom:
            break
        
        # "-" ile başlayan satırlar için GRİ ARKA PLANLI KUTU
        if line.startswith("-") or line.startswith("−"):
            text_content = line[1:].strip()
            
            # Metin uzunluğuna göre kutu yüksekliği hesapla
            # Yaklaşık 90 karakter bir satıra sığar (font 16, width 12 inch)
            chars_per_line = 85
            num_lines = max(1, (len(text_content) + chars_per_line - 1) // chars_per_line)
            dynamic_box_height = 0.35 + (num_lines * 0.25)  # Her satır için 0.25 inch ekle
            
            # Gri arka planlı rounded rectangle
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, 
                left, current_top,
                width, Inches(dynamic_box_height)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = GRI_KUTU
            shape.line.fill.background()
            # Köşe yuvarlaklığı
            shape.adjustments[0] = 0.08
            
            # Metin ekle
            text_box = slide.shapes.add_textbox(
                left + Inches(0.15), 
                current_top + Inches(0.1), 
                width - Inches(0.3), 
                Inches(dynamic_box_height - 0.15)
            )
            tf = text_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text_content
            p.font.size = Pt(font_size)
            p.font.color.rgb = BEYAZ
            p.font.name = "Arial"
            current_top += Inches(dynamic_box_height + 0.08)
        
        # "•" ile başlayan maddeler - HER ZAMAN MADDE İŞARETİ
        elif line.startswith("•"):
            text_content = line[1:].strip()
            text_box = slide.shapes.add_textbox(
                left + Inches(0.3), current_top, 
                width - Inches(0.3), Inches(0.35)
            )
            tf = text_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "●  " + text_content
            p.font.size = Pt(font_size - 1)
            p.font.color.rgb = ACIK_GRI
            p.font.name = "Arial"
            current_top += Inches(line_height)
        
        # "//" ile başlayan alt başlıklar
        elif line.startswith("//"):
            current_top += Inches(0.08)
            text_box = slide.shapes.add_textbox(left, current_top, width, Inches(0.4))
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = line[2:].strip()
            p.font.size = Pt(font_size + 1)
            p.font.bold = True
            p.font.color.rgb = ACCENT
            p.font.name = "Arial"
            current_top += Inches(line_height + 0.05)
        
        # Normal metin
        else:
            text_box = slide.shapes.add_textbox(left, current_top, width, Inches(0.4))
            tf = text_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = line
            p.font.size = Pt(font_size)
            p.font.color.rgb = ACIK_GRI
            p.font.name = "Arial"
            current_top += Inches(line_height)
    
    return current_top

def add_slide_number(slide, num, total):
    """Slayt numarası ekle"""
    left, top = Inches(12.5), Inches(7.0)
    text_box = slide.shapes.add_textbox(left, top, Inches(0.7), Inches(0.3))
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.RIGHT

# ============ SLAYTLAR ============

slides_data = [
    # Slayt 1: Kapak
    {
        "title": "OTTOBITE",
        "subtitle": "Garson Davranışları ve İş Önceliği Rehberi",
        "is_cover": True
    },
    # Slayt 2: Amaç
    {
        "title": "Amaç",
        "content": [
            "- Bu doküman, OTTOBITE salon servisinde görev alan tüm servis personelinin;",
            "• Doğru davranış biçimini benimsemesi,",
            "• Yoğun servis anlarında doğru iş önceliğini kurabilmesi,",
            "• Salon düzenini ve misafir memnuniyetini koruyabilmesi",
            "amacıyla hazırlanmıştır.",
            "",
            "- Bu rehber okunmak için değil, uygulanmak için hazırlanmıştır."
        ]
    },
    # Slayt 3: 1. Garsonluk Bilinci ve Temel Zihniyet
    {
        "title": "1. Garsonluk Bilinci ve Temel Zihniyet",
        "content": [
            "- Garsonluk yalnızca sipariş almak değildir.",
            "- Garson, salonun temposunu, misafirin ruh halini ve servisin akışını yönetir.",
            "",
            "Misafir, işletmeyi;",
            "• Garsonun bakışından,",
            "• Konuşma tarzından,",
            "• Salondaki duruşundan",
            "anlar ve değerlendirir.",
            "",
            "OTTOBITE'de garson:",
            "• Göz önündedir",
            "• Dikkatlidir",
            "• Sorumluluk sahibidir",
            "",
            "- Şu unutulmamalıdır: Garsonun davranışı, işletmenin kalitesini temsil eder."
        ]
    },
    # Slayt 4: 2. Temel Garson Davranış Kuralları
    {
        "title": "2. Temel Garson Davranış Kuralları",
        "content": [
            "- Bu bölümde yer alan kurallar tartışmaya açık değildir.",
            "- Tüm servis personeli tarafından eksiksiz uygulanır.",
            "",
            "// Duruş ve Salon İçindeki Tavır",
            "• Garson salonda dik ve hazır durur",
            "• El cepte gezilmez",
            "• Kollar bağlı şekilde durulmaz",
            "• Duvara yaslanılmaz",
            "• Boş boş ayakta beklenmez",
            "",
            "Garson salondayken: \"Her an çağrılabilirim\" bilinciyle hareket eder."
        ]
    },
    # Slayt 5: 2. Temel Garson Davranış Kuralları (devam)
    {
        "title": "2. Temel Garson Davranış Kuralları (devam)",
        "content": [
            "// Masaya Yaklaşma ve İletişim",
            "• Masaya yaklaşırken göz teması kurulur",
            "• Masaya sırt dönülerek konuşulmaz",
            "• Misafirin yüzüne bakmadan sipariş alınmaz",
            "",
            "Kullanılmaması gereken ifadeler:",
            "• \"Bir saniye\"",
            "• \"Dur bakayım\"",
            "• \"Yoğunum\"",
            "",
            "Kullanılması gereken yaklaşım:",
            "• \"Hemen ilgileniyorum\"",
            "• \"Şimdi yardımcı oluyorum\"",
            "",
            "- Ton sakin, net ve kontrollü olmalıdır."
        ]
    },
    # Slayt 6: 2. Temel Garson Davranış Kuralları (devam 2)
    {
        "title": "2. Temel Garson Davranış Kuralları (devam)",
        "content": [
            "// Salon İçinde Davranış Disiplini",
            "• Personel kendi arasında yüksek sesle konuşmaz",
            "• Gülüşmeler misafir dikkatini çekecek seviyede olmaz",
            "• Telefon kesinlikle kullanılmaz",
            "",
            "Şef salondayken:",
            "• Garson görünür olur",
            "• Kaçmaz",
            "• Saklanmaz"
        ]
    },
    # Slayt 7: 3. Yoğun Serviste Garsonun Düşünme Şekli
    {
        "title": "3. Yoğun Serviste Garsonun Düşünme Şekli",
        "content": [
            "- Bu bölüm, servis kalitesini doğrudan belirler.",
            "- Yoğunlukta hızlı olmak değil, doğru öncelik vermek önemlidir.",
            "- Panik, hatayı artırır.",
            "- Öncelik, düzeni sağlar.",
            "",
            "// Doğru İş Önceliği Sıralaması",
            "OTTOBITE servisinde garson şu sırayı bilmelidir:",
            "• Yeni gelen misafir / boş masa",
            "• Sipariş vermek isteyen masa",
            "• Servisi çıkan masa",
            "• Hesap isteyen masa",
            "• Diğer işler",
            "",
            "Bu sıralama yoğunlukta değişmez."
        ]
    },
    # Slayt 8: 3. Yoğun Serviste (devam)
    {
        "title": "3. Yoğun Serviste Garsonun Düşünme Şekli (devam)",
        "content": [
            "Her servis türü;",
            "• Dikkat,",
            "• Düzen,",
            "• Misafirle doğru iletişim",
            "gerektirir.",
            "",
            "- Garson servis türüne göre zihinsel olarak hazırlanmalı ve her servisi aynı profesyonellikle yürütmelidir.",
            "",
            "// Aynı Anda Her Şey Yapılmaz",
            "• Aynı anda 5 masaya yetişmeye çalışılmaz",
            "• Bir işi yarım bırakıp diğerine geçilmez",
            "• Yapılan iş tamamlanır, sonra diğerine geçilir",
            "",
            "- Şef yönlendirmesi her zaman önceliklidir."
        ]
    },
    # Slayt 9: 4. Yoğun Serviste Yapılmaması Gereken Davranışlar
    {
        "title": "4. Yoğun Serviste Yapılmaması Gereken Davranışlar",
        "content": [
            "Aşağıdaki davranışlar servis kalitesini düşürür ve kabul edilmez:",
            "",
            "• \"Ben o masaya bakmıyorum\" demek",
            "• Masayı görmezden gelmek",
            "• Mutfaktan veya salondan kaçmak",
            "• Şefe savunma yapmak",
            "• Misafir önünde tartışmak",
            "",
            "- Bu davranışlar OTTOBITE servis anlayışına uygun değildir."
        ]
    },
    # Slayt 10: 5. Garson – Salon Şefi İlişkisi
    {
        "title": "5. Garson – Salon Şefi İlişkisi",
        "content": [
            "Salon şefi, servisin düzeninden sorumludur.",
            "",
            "Garson şunu bilmelidir:",
            "• Şef uyarısı kişisel değildir",
            "• Şef yönlendirmesi tartışılmaz",
            "• Servis sırasında savunma yapılmaz",
            "",
            "- Servis anında uygulama yapılır, değerlendirme servis sonrası yapılır."
        ]
    },
    # Slayt 11: 7. Açılış ve Kapanış Disiplini
    {
        "title": "7. Açılış ve Kapanış Disiplini",
        "content": [
            "Açılış ve kapanış saatlerinde uyuşukluk kabul edilemez.",
            "",
            "• Yavaş hareket etmek",
            "• Eğlence modunda çalışmak",
            "• Sürekli sohbet halinde olmak",
            "• 0.5x hızla iş yapmak",
            "",
            "kabul edilemez davranışlardır.",
            "",
            "Açılış ve kapanış:",
            "• Dikkat ister",
            "• Odak ister",
            "• Disiplin ister",
            "",
            "- Açılış ve kapanış ne kadar düzgün yapılırsa, vardiya o kadar rahat geçer."
        ]
    },
    # Slayt 12: 8. İş Hızı ve İş Önceliği Bilinci
    {
        "title": "8. İş Hızı ve İş Önceliği Bilinci",
        "content": [
            "Bir garson şunu bilmelidir:",
            "• 3 dakikalık bir iş 15 dakikaya uzatılamaz",
            "• Yavaşlık beceri değildir",
            "• Odağın dağılması iş aksatır",
            "",
            "Kendi postasında yapacağı 4 basit işi;",
            "• Sohbet ederek",
            "• Etrafla ilgilenerek",
            "• Önceliği yanlış belirleyerek",
            "uzatmak kabul edilemez.",
            "",
            "- İşte olduğunu bilmek ve buna göre davranmak zorunluluktur."
        ]
    },
    # Slayt 13: 9. İş Esnasında Tutum ve Beden Dili
    {
        "title": "9. İş Esnasında Tutum ve Beden Dili",
        "content": [
            "Servis sırasında:",
            "• Sürekli şikâyet etmek",
            "• Yakınmak",
            "• Uflamak, puflamak",
            "• Yüz mimikleriyle ima yapmak",
            "",
            "kesinlikle kabul edilemez.",
            "",
            "Garsonun yüz ifadesi ve beden dili;",
            "• Misafire",
            "• Takım arkadaşlarına",
            "• Operasyona",
            "olumsuz yansır."
        ]
    },
    # Slayt 14: 10. Gruplaşma, Dedikodu ve İletişim Sorunları
    {
        "title": "10. Gruplaşma, Dedikodu ve İletişim Sorunları",
        "content": [
            "• Gruplaşma",
            "• Dedikodu",
            "• Kulis yapmak",
            "",
            "OTTOBITE çalışma düzenine aykırıdır.",
            "",
            "Farklı departman şefleriyle veya personelleriyle:",
            "• Tartışmaya girmek",
            "• Münakaşa yaşamak",
            "",
            "kabul edilemez.",
            "",
            "- Sorunlar servis dışında ve doğru kanallarla iletilir."
        ]
    },
    # Slayt 15: 11. Kurallara Uyum ve Profesyonel Davranış
    {
        "title": "11. Kurallara Uyum ve Profesyonel Davranış",
        "content": [
            "Yönetim tarafından belirlenen kuralları:",
            "• Esnetmek",
            "• Suistimal etmek",
            "• Kendi lehine yorumlamak",
            "",
            "kabul edilemez.",
            "",
            "Herkes:",
            "• İş esnasında uyumlu olmak",
            "• Kendi görevini en iyi şekilde yapmak",
            "zorundadır."
        ]
    },
    # Slayt 16: 12. Molalar ve Görev Takibi
    {
        "title": "12. Molalar ve Görev Takibi",
        "content": [
            "• Hiçbir personel başka bir personelin molasını takip etmez",
            "• Molalar yönetim tarafından düzenlenir",
            "",
            "- Bu konuda bireysel kontrol veya yorum yapılmaz."
        ]
    },
    # Slayt 17: 13. Posta Yönetimi ve Masa Kontrolü
    {
        "title": "13. Posta Yönetimi ve Masa Kontrolü",
        "content": [
            "Garson kendi postasındaki masalardan tamamen sorumludur.",
            "",
            "Sürekli uyarılan durumlar:",
            "• Toplanmamış peçeteler",
            "• Islak mendil çöpleri",
            "• Şişe boşları",
            "• Değiştirilmemiş kül tablaları",
            "• Masa üzerindeki müşteri çöpleri",
            "",
            "- İyi bir garson kendi postasındaki 5 masayı eksiksiz yönetir.",
            "- Sürekli uyarı gerektiren çalışma kabul edilemez."
        ]
    },
    # Slayt 18: 14. Kişisel Görünüm ve Kıyafet Disiplini
    {
        "title": "14. Kişisel Görünüm ve Kıyafet Disiplini",
        "content": [
            "• Bayan personel saçını açık kullanamaz",
            "• İş kıyafeti ne ise o giyilir",
            "• Kıyafet dışında kişisel yorum yapılmaz",
            "",
            "- Görünüm işletmenin ciddiyetini yansıtır."
        ]
    },
    # Slayt 19: 15. İletişim Bilinci (EN ÖNEMLİ KONU)
    {
        "title": "15. İletişim Bilinci (EN ÖNEMLİ KONU)",
        "content": [
            "OTTOBITE servisinin en önemli unsuru iletişimdir.",
            "",
            "• Servisin hatasız ilerlemesi",
            "• Posta kontrolü",
            "• Yoğunluk yönetimi",
            "",
            "ancak doğru iletişimle sağlanır.",
            "",
            "Posta terk edilecekse:",
            "• Takım arkadaşları bilgilendirilir",
            "• Şef haberdar edilir",
            "",
            "- Habersiz posta terk edilemez."
        ]
    },
    # Slayt 20: 16. Hiyerarşi ve Saygı
    {
        "title": "16. Hiyerarşi ve Saygı",
        "content": [
            "• Hiçbir personel başka bir personele üstü gibi davranamaz",
            "• Yetki yalnızca görev tanımıyla sınırlıdır",
            "",
            "- Saygı karşılıklıdır ve herkes için geçerlidir."
        ]
    },
    # Slayt 21: 17. İşe Zamanında Gelme Disiplini
    {
        "title": "17. İşe Zamanında Gelme Disiplini",
        "content": [
            "İşe geç kalmak;",
            "• Operasyonu aksatır",
            "• Takım arkadaşlarını zor durumda bırakır",
            "• Servis düzenini bozar",
            "",
            "- Geç kalmak bir alışkanlık hâline gelemez.",
            "",
            "Her personel vardiyasına:",
            "• Zamanında",
            "• Hazır",
            "• İşe odaklanmış şekilde",
            "başlamak zorundadır."
        ]
    },
    # Slayt 22: 18. Masa Kalkış Sonrası Temizlik Bilinci
    {
        "title": "18. Masa Kalkış Sonrası Temizlik Bilinci",
        "content": [
            "Misafir kalktığında yalnızca masa üstü değil;",
            "• Masa altı",
            "• Sandalyeler",
            "• Genel düzen",
            "kontrol edilir.",
            "",
            "- Masa altına faraş atılmadan yeni misafir alınmaz.",
            "- Temizlenmeyen masa, yarım yapılmış iştir."
        ]
    },
    # Slayt 23: 19. Yardımlaşma ve Takım Bilinci
    {
        "title": "19. Yardımlaşma ve Takım Bilinci",
        "content": [
            "- Servis bireysel değil, ekip işidir.",
            "",
            "• Kendi postasında işi olan garson çalışır",
            "• Diğer postadaki garson müsaitse destek verir",
            "",
            "Eğer:",
            "- Senin aklında 5 iş varsa diğer posta boş duruyorsa yardım istemek zorunluluktur."
        ]
    },
    # Slayt 24: 20. İletişimde Esneklik ve Açıklık
    {
        "title": "20. İletişimde Esneklik ve Açıklık",
        "content": [
            "Servisin aksamaması için:",
            "• Herkes birbirine konuşabilmelidir",
            "• Kimse bir şey söylemeye çekinmemelidir",
            "",
            "- Bir personel, başka bir personele bir şey demeye çekinirse iş aksar.",
            "",
            "İletişimde:",
            "• Esneklik",
            "• Açıklık",
            "• Netlik",
            "zorunludur."
        ]
    },
    # Slayt 25: 21. Sessiz Tepki ve İma Yasağı
    {
        "title": "21. Sessiz Tepki ve İma Yasağı",
        "content": [
            "• İmalı konuşmalar",
            "• Sessiz tepkiler",
            "• Yüz mimikleriyle mesaj verme",
            "",
            "kabul edilemez.",
            "",
            "- Sorun varsa açıkça ve doğru yerde dile getirilir."
        ]
    },
    # Slayt 26: 22. Boş Durma ve Oyalama Yasağı
    {
        "title": "22. Boş Durma ve Oyalama Yasağı",
        "content": [
            "Servis sırasında:",
            "• Boş boş gezmek",
            "• Oyalandığı izlenimi vermek",
            "• Gereksiz hareketlerle zaman harcamak",
            "",
            "kabul edilemez.",
            "",
            "- Her an yapılacak bir iş vardır."
        ]
    },
    # Slayt 27: 23. Alan Sahiplenme Bilinci
    {
        "title": "23. Alan Sahiplenme Bilinci",
        "content": [
            "- Her garson kendi alanını sahiplenir.",
            "",
            "• \"Orası benim alanım değil\" yaklaşımı yoktur",
            "• Salon herkesindir",
            "",
            "- İhtiyaç olan yerde destek olunur."
        ]
    },
    # Slayt 28: 24. Mutfak ve Bar ile İletişim
    {
        "title": "24. Mutfak ve Bar ile İletişim",
        "content": [
            "Mutfak ve bar:",
            "• Tartışma alanı değildir",
            "• Emir verilen yer değildir",
            "",
            "- İletişim saygılı ve net olmalıdır."
        ]
    },
    # Slayt 29: 25. Yoğunlukta Sessiz ve Kontrollü Çalışma
    {
        "title": "25. Yoğunlukta Sessiz ve Kontrollü Çalışma",
        "content": [
            "Yoğunluk:",
            "• Ses yükseltme",
            "• Koşturma",
            "• Panik",
            "sebebi değildir.",
            "",
            "- Kontrollü çalışan garson, yoğunluğu yönetir."
        ]
    },
    # Slayt 30: 26. Misafir Önünde Ekip Disiplini
    {
        "title": "26. Misafir Önünde Ekip Disiplini",
        "content": [
            "Misafir önünde:",
            "• Uyarı yapılmaz",
            "• Tartışılmaz",
            "• Rol kesilmez",
            "",
            "- Tüm ekip tek vücut gibi görünmelidir."
        ]
    },
    # Slayt 31: 27. Kişisel Sorunları İşe Yansıtmama
    {
        "title": "27. Kişisel Sorunları İşe Yansıtmama",
        "content": [
            "Kişisel problemler:",
            "• Yüz ifadelerine",
            "• Tavırlara",
            "• Konuşmalara",
            "yansıtılamaz."
        ]
    },
    # Slayt 32: 28. İş Bitmeden Rahatlama Yasağı
    {
        "title": "28. İş Bitmeden Rahatlama Yasağı",
        "content": [
            "İş bitmeden:",
            "• Rahatlama",
            "• Dağılma",
            "• Tempo düşürme",
            "kabul edilemez."
        ]
    },
    # Slayt 33: 29. Tekrar Eden Hatalar
    {
        "title": "29. Tekrar Eden Hatalar",
        "content": [
            "Aynı hatanın sürekli yapılması:",
            "• Dikkatsizliktir",
            "• Umursamazlıktır",
            "",
            "Kabul edilemez."
        ]
    },
    # Slayt 34: 30. Kendi Kendini Kontrol Etme
    {
        "title": "30. Kendi Kendini Kontrol Etme",
        "content": [
            "İyi bir garson:",
            "• Uyarı beklemez",
            "• Kendini kontrol eder",
            "• Eksiklerini fark eder"
        ]
    },
    # Slayt 35: 31. Vardiya Boyunca Hazır Olma
    {
        "title": "31. Vardiya Boyunca Hazır Olma",
        "content": [
            "Vardiya boyunca:",
            "• Zihinsel",
            "• Fiziksel",
            "olarak hazır olmak zorunludur."
        ]
    },
    # Slayt 36: 32. Gereksiz Yorum ve Mizah Yasağı
    {
        "title": "32. Gereksiz Yorum ve Mizah Yasağı",
        "content": [
            "Servis sırasında:",
            "• Gereksiz şaka",
            "• Aşırı samimiyet",
            "kabul edilemez."
        ]
    },
    # Slayt 37: 33. Görev Tanımına Sadakat
    {
        "title": "33. Görev Tanımına Sadakat",
        "content": [
            "- Herkes kendi görev tanımı içinde hareket eder.",
            "",
            "- Yetki aşımı yapılamaz."
        ]
    },
    # Slayt 38: 34. Ortak Amaç Bilinci
    {
        "title": "34. Ortak Amaç Bilinci",
        "content": [
            "Herkesin ortak amacı:",
            "• Düzgün servis",
            "• Hatasız operasyon",
            "• Memnun misafir",
            "olmalıdır."
        ]
    },
    # Slayt 39: 35. QR Değerlendirme ve Misafir Geri Bildirimi
    {
        "title": "35. QR Değerlendirme ve Misafir Geri Bildirimi",
        "content": [
            "Servis kalitesinin ölçülmesi için:",
            "• İsim verilerek selamlanmalıdır",
            "• Servis sahiplenilmelidir",
            "• Misafir görmezden gelinemez",
            "",
            "- \"Merhaba, ben Ahmet. Servisinizle ben ilgileneceğim.\"",
            "",
            "// QR Talebi ve Zamanlama",
            "• Servis bitimi beklenir",
            "• Sohbet bölünmez",
            "• Baskı kurulmaz"
        ]
    },
    # Slayt 40: 35. QR Değerlendirme (devam)
    {
        "title": "35. QR Değerlendirme ve Misafir Geri Bildirimi (devam)",
        "content": [
            "- \"Sizleri ağırlamak bir keyifti. Hizmet kalitemiz hakkındaki olumlu geri bildirimleriniz, ekibimiz için motivasyon kaynağı olacaktır.\"",
            "",
            "// İletişim ve Tavır",
            "• Göz teması",
            "• Nötr ve nazik yüz ifadesi",
            "• Profesyonel mesafe",
            "zorunludur.",
            "",
            "- Savunma, ima ve gereksiz mizah kabul edilemez.",
            "- Her garson, servisinin puanlanacağı bilinciyle hareket eder."
        ]
    },
    # Slayt 41: 36. Memnuniyetsizlik Yönetimi ve Çözüm Süreci
    {
        "title": "36. Memnuniyetsizlik Yönetimi ve Çözüm Süreci",
        "content": [
            "Misafir şikayeti veya olumsuz geri bildirim durumunda:",
            "• İlgili yöneticiye anında haber verilir",
            "• Personel kendi başına karar almaz",
            "• Kişisel çözüm sürecine girilmez",
            "",
            "- Yönetici bilgilendirmesi zorunludur.",
            "",
            "// İnisiyatif ve Esneklik",
            "• Misafirin mutlu ayrılması esastır",
            "• Operasyonel kurallar yönetici kontrolünde esnetilebilir"
        ]
    },
    # Slayt 42: 36. Memnuniyetsizlik (devam)
    {
        "title": "36. Memnuniyetsizlik Yönetimi ve Çözüm Süreci (devam)",
        "content": [
            "// Bilgilendirme Yükümlülüğü",
            "- Sunulan çözüm, \"istisnai bir durum\" olduğu vurgulanarak misafire aktarılır.",
            "",
            "// Örnek Diyalog",
            "- \"Normalde kahvaltı servisimiz 13:00'te sona ermektedir; ancak bugün mutfağımızla görüşerek size özel bir istisna sağladık. Keyifli bir deneyim dileriz.\"",
            "",
            "- Kendi başına kural esnetmek veya sorunu gizlemek kabul edilemez."
        ]
    },
    # Slayt 43: 37. Koordineli Mola ve İş Akışı Yönetimi
    {
        "title": "37. Koordineli Mola ve İş Akışı Yönetimi",
        "content": [
            "Servis akışının kesintisiz ilerlemesi için:",
            "• Birbiriyle bağlantılı birimler aynı anda molaya çıkamaz",
            "• Mola planlaması, operasyonel boşluk yaratmayacak şekilde yapılır",
            "• Ekip içi koordinasyon mola öncesinde tamamlanır",
            "",
            "- Hizmetin sürekliliği esastır.",
            "",
            "// Birim Çakışması",
            "• Bar / Bar",
            "• Servis / Servis",
            "• Kasa / Vitrin"
        ]
    },
    # Slayt 44: 37. Koordineli Mola (devam)
    {
        "title": "37. Koordineli Mola ve İş Akışı Yönetimi (devam)",
        "content": [
            "- Aynı görev tanımına sahip veya birbirini yedeklemesi gereken birimlerin aynı anda mola planlaması kabul edilemez.",
            "",
            "// Kritik Kural",
            "- Mola saatlerinde bireysel istek değil, operasyonel ihtiyaç önceliklidir.",
            "- Yönetici onayı ve birimler arası mutabakat zorunludur."
        ]
    },
    # Slayt 45: 38. Boş Toplama ve Tepsi Düzeni
    {
        "title": "38. Boş Toplama ve Tepsi Düzeni",
        "content": [
            "Servis güvenliği ve operasyonel hız için:",
            "• Ağırlık merkezde toplanmalıdır",
            "• Büyükten küçüğe doğru istifleme yapılmalıdır",
            "• Kırılma riski olan parçalar ayrıştırılmalıdır",
            "",
            "// Tepsi Dengesi",
            "• Ağır ve büyük parçalar merkeze yerleştirilir",
            "• Bardaklar ve tabaklar kontrolsüz şekilde yığılamaz"
        ]
    },
    # Slayt 46: 38. Boş Toplama (devam)
    {
        "title": "38. Boş Toplama ve Tepsi Düzeni (devam)",
        "content": [
            "// Hatalı İstifleme",
            "• Dengesi bozuk yükleme",
            "• Görüşü engelleyecek yükseklikte yığma",
            "• Tabakların kaymasına neden olacak dizim",
            "kabul edilemez.",
            "",
            "- Düzenli dizilen tepsi; profesyonel görünüm ve iş güvenliği için zorunludur."
        ]
    },
    # Slayt 47: 39. Misafir Karşılama ve Proaktif Takip
    {
        "title": "39. Misafir Karşılama ve Proaktif Takip",
        "content": [
            "Salonun canlılığı ve misafir konforu için:",
            "• Kapıdan giren her misafir anında fark edilmelidir",
            "• Göz teması ve tebessüm ile \"Hoş geldiniz\" karşılaması yapılır",
            "• Misafir, masaya yerleşene kadar eşlik edilir",
            "",
            "- Müşteri kapıda bekletilemez, görmezden gelinemez."
        ]
    },
    # Slayt 48: 39. Misafir Karşılama (devam)
    {
        "title": "39. Misafir Karşılama ve Proaktif Takip (devam)",
        "content": [
            "// İhtiyaç Hissetme (Proaktif Hizmet)",
            "• Misafir garsonu aramaz; garson misafiri izler",
            "• Biten içecekler, azalan peçeteler ve boşalan tabaklar takip edilir",
            "",
            "- Misafirin çevreye bakması veya el kaldırması \"hizmet gecikmesi\" olarak kabul edilir.",
            "",
            "// Takip Esasları",
            "• Boşalan bardağın yenilenmesi teklif edilir",
            "• İhtiyaç duyulmadan peçete/servis desteği sağlanır",
            "• Misafirle göz teması kurmak için salon sürekli taranır"
        ]
    },
    # Slayt 49: 39. Misafir Karşılama (devam 2)
    {
        "title": "39. Misafir Karşılama ve Proaktif Takip (devam)",
        "content": [
            "- Misafiri el kaldırmaya veya seslenmeye mecbur bırakmak kabul edilemez."
        ]
    },
    # Slayt 50: 40. Hazırlık ve Arka Plan İşleri Yönetimi
    {
        "title": "40. Hazırlık ve Arka Plan (Back) İşleri Yönetimi",
        "content": [
            "Servis kalitesinin sürekliliği için:",
            "• Çatal-bıçak silme",
            "• Takım ve ekipman çekme",
            "• İstifleme ve düzenleme",
            "işleri, servis yoğunluğu ve misafir trafiğine göre planlanır.",
            "",
            "// Öncelik Sıralaması",
            "• Misafir her zaman hazırlık işlerinden önceliklidir",
            "• Salon boş bırakılarak arka plana odaklanılamaz",
            "• Yoğunluk anında büyük çaplı hazırlık işlerine girilmesi kabul edilemez"
        ]
    },
    # Slayt 51: 40. Hazırlık (devam)
    {
        "title": "40. Hazırlık ve Arka Plan (Back) İşleri Yönetimi (devam)",
        "content": [
            "// Zamanlama ve Uygulama",
            "• Hazırlıklar servis öncesi veya sakin saatlerde tamamlanır",
            "• Mutfaktan takım çekme ve silme işlemleri misafiri rahatsız etmeyecek şekilde sessizce yürütülür",
            "• Eksikler, servis krize girmeden proaktif (önceden) olarak giderilir",
            "",
            "- Servis önceliği esastır.",
            "- Arka plan işleri, servis akışını engellemek için değil, desteklemek içindir."
        ]
    },
    # Slayt 52: 41. Operasyonel İletişimde Netlik ve Kısalık
    {
        "title": "41. Operasyonel İletişimde Netlik ve Kısalık",
        "content": [
            "Hızlı ve hatasız iş akışı için:",
            "• İletişim en az kelimeyle, en net şekilde kurulmalıdır",
            "• Uzun açıklamalar yerine sonuç odaklı cümleler seçilir",
            "• Mesaj, karşı tarafın işini aksatmayacak hızda aktarılır",
            "",
            "- \"3 kelimeyle anlatılabilecek bir durum, gereksiz detaylarla uzatılamaz.\"",
            "",
            "// Uzun ve Karmaşık Anlatım",
            "• Dinleyicinin odağını dağıtır",
            "• Diğer işlerin unutulmasına neden olur",
            "• Operasyonu yavaşlatır"
        ]
    },
    # Slayt 53: 41. Operasyonel İletişim (devam)
    {
        "title": "41. Operasyonel İletişimde Netlik ve Kısalık (devam)",
        "content": [
            "// Anlatım Bozukluğu ve Karmaşa",
            "• Olayı saptırmak",
            "• Konuyu dolandırmak",
            "• Gereksiz detaylara girmek",
            "kabul edilemez.",
            "",
            "- Hızlı, doğrudan ve net iletişim zorunludur."
        ]
    },
    # Slayt 54: 42. Son Hatırlatma
    {
        "title": "42. Son Hatırlatma",
        "content": [
            "OTTOBITE'de çalışmak;",
            "• Disiplin",
            "• Dikkat",
            "• Takım bilinci",
            "ister.",
            "",
            "- Bu rehberde yer alan tüm maddeler uygulanmak zorundadır."
        ]
    },
    # Slayt 55: Ana Fikir
    {
        "title": "Ana Fikir",
        "content": [
            "- OTTOBITE çatısı altında hangi birimde olursak olalım, her birimiz markamızın yaşayan birer temsilcisiyiz.",
            "",
            "- Temel değerimiz; birimler arası kusursuz bir uyum ve yüksek sorumluluk bilinciyle hareket ederek, misafirlerimize yalnızca bir hizmet değil, disiplin ve nezaketle örülmüş bir deneyim sunmaktır.",
            "",
            "- Bizim için başarı; mutfaktan kasaya, servisten bara kadar her noktada aynı profesyonel tavrı sergileyen \"tek bir ekip\" olabilmektir."
        ],
        "is_closing": True
    }
]

total_slides = len(slides_data)

for idx, slide_data in enumerate(slides_data):
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Gradient background
    add_gradient_background(slide)
    
    if slide_data.get("is_cover"):
        # Kapak slaytı
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data["title"]
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = KIREMIT
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER
        
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data["subtitle"]
        p.font.size = Pt(28)
        p.font.color.rgb = ACIK_GRI
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER
        
    elif slide_data.get("is_closing"):
        # Kapanış slaytı
        add_title_shape(slide, slide_data["title"], is_main=True)
        add_content_box(slide, slide_data["content"], start_top=1.5, base_font_size=17)
    else:
        # Normal içerik slaytları
        add_title_shape(slide, slide_data["title"])
        add_content_box(slide, slide_data["content"])
    
    # Slayt numarası
    add_slide_number(slide, idx + 1, total_slides)

# Sunum kaydet
output_path = "/Users/ahmet11159/Desktop/OTTOBITE_Garson_Rehberi.pptx"
prs.save(output_path)
print(f"Sunum başarıyla oluşturuldu: {output_path}")
print(f"Toplam slayt sayısı: {total_slides}")
