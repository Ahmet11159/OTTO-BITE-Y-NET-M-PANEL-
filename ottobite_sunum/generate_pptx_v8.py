# encoding: utf-8
import json
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR

def create_presentation():
    prs = Presentation()
    BRICK_RED = RGBColor(192, 57, 43)
    DARK_GRAY = RGBColor(44, 62, 80)
    LIGHT_BG = RGBColor(245, 246, 250) # Very light gray for boxes
    
    with open('structured_data_final_v6.json', 'r', encoding='utf-8') as f:
        slides_data = json.load(f)

    # --- Auto-Fit Helper ---
    def fit_text_content(tf, content_items, max_height_inches=5.8):
        """
        Populate text frame and iteratively turn down font size if it doesn't fit?
        Actually, python-pptx doesn't give calculated height easily.
        We must use strict heuristic calculation.
        """
        
        # Try sizes: 20, 18, 16, 14, 12, 10
        sizes = [20, 18, 16, 14, 12, 11, 10]
        
        selected_size = 12 # Default fallback
        
        for size in sizes:
            current_height = 0
            # Constants per size
            line_height = size * 1.4 # points
            para_spacing = size * 0.6 # points
            
            total_points = 0
            
            for item in content_items:
                text = item['text']
                # Wrap estimation: Avg char width approx 0.5 * size
                # Textbox width approx 9 inches = 648 points
                # Max chars per line approx 648 / (0.5 * size)
                max_chars = 900 / (0.6 * size) # rough est
                
                lines = math.ceil(len(text) / max_chars)
                if lines < 1: lines = 1
                
                # Height for this item
                item_h = (lines * line_height) + para_spacing
                if item['type'] == 'subheader': item_h += 10
                if item['type'] == 'quote_box': item_h += 15
                
                total_points += item_h
            
            # Convert total points to inches
            total_inches = total_points / 72
            if total_inches <= max_height_inches:
                selected_size = size
                break
        
        # Apply selected size
        for item in content_items:
            text_val = item['text']
            item_type = item['type']
            
            p = tf.add_paragraph()
            
            # STYLING
            if item_type == "subheader":
                p.text = text_val
                p.font.bold = True
                p.font.size = Pt(selected_size + 2)
                p.font.color.rgb = DARK_GRAY
                p.space_before = Pt(selected_size)
                
            elif item_type == "bullet":
                p.text = "• " + text_val
                p.font.size = Pt(selected_size)
                p.font.color.rgb = DARK_GRAY
                p.space_after = Pt(selected_size * 0.4)
                
            elif item_type == "emphasis":
                p.text = text_val.upper()
                p.font.bold = True
                p.font.color.rgb = RGBColor(192, 0, 0)
                p.font.size = Pt(selected_size)
                p.space_before = Pt(selected_size * 0.5)
                p.alignment = PP_ALIGN.CENTER
                
            elif item_type == "quote_box" or item_type == "body_text":
                # Simulate a box using a separate shape? No, mixed content in one textbox is hard.
                # Just style it distinctively: Italic, different color background?
                # Textframe paragraph background not supported easily.
                # We will use Indent + Borders? No.
                # Just use distinct text style: Darker, Italic, possibly centered or indented.
                
                p.text = text_val
                p.font.size = Pt(selected_size)
                p.font.italic = True
                if item_type == "quote_box":
                    p.font.color.rgb = RGBColor(60, 60, 60)
                    p.level = 0
                    # Add simple "quote" visual indicator
                    p.text = "“ " + text_val.replace('“','').replace('"','') + " ”"
                else:
                    p.font.color.rgb = DARK_GRAY
                
                p.space_after = Pt(selected_size * 0.5)
            
            p.font.name = "Calibri"

    # --- Title Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    p = txBox.text_frame.add_paragraph()
    p.text = "OTTOBITE"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = BRICK_RED
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1))
    p2 = txBox2.text_frame.add_paragraph()
    p2.text = "Garson Davranışları ve İş Önceliği Rehberi"
    p2.font.size = Pt(28)
    p2.font.color.rgb = DARK_GRAY
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = "Calibri Light"
    
    # --- Content Slides ---
    for slide_info in slides_data:
        if len(slide_info['content']) == 0: continue
        if slide_info['title'] == "OTTOBITE Garson Rehberi" and len(slide_info['content']) < 3: continue
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 1. Header with Dynamic Line
        title_text = slide_info['title']
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
        tp = title_box.text_frame.add_paragraph()
        tp.text = title_text
        tp.font.size = Pt(32)
        tp.font.bold = True
        tp.font.color.rgb = BRICK_RED
        tp.font.name = "Arial"
        
        # Estimate Line Width based on char count
        # Avg char width for Arial Bold 32pt is approx 0.18 - 0.22 inches? 
        # Let's say 0.23 inches per char.
        char_count = len(title_text)
        est_width = char_count * 0.23
        if est_width > 9: est_width = 9
        if est_width < 1: est_width = 1
        
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), est_width, Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = BRICK_RED
        line.line.fill.background()
        
        # 2. Decorate "Box Text" (Intro/Quotes)
        # Scan for box_text items to create separate background shapes BEHIND the text?
        # Too complex positioning.
        # Alternative: Place a grey rounded rectangle for specific items if found?
        # Let's stick to the Auto-Fit Logic inside one main text box for alignment safety.
        # But for "Visual Professionalism", the user requested a "window".
        
        # Hybrid Approach: 
        # Identify "Intro" content (Body Text at start).
        intro_items = []
        main_items = []
        
        for item in slide_info['content']:
            if item['type'] in ['body_text', 'quote_box'] and not main_items:
                intro_items.append(item)
            else:
                main_items.append(item)
        
        top_cursor = 1.6
        
        # Render Intro Box if exists
        if intro_items:
            # Create a shape for intro
            # Estimate height?
            est_lines = sum([math.ceil(len(i['text'])/80) for i in intro_items])
            est_height = est_lines * 0.4 + 0.3
            
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, 
                Inches(0.5), Inches(top_cursor), Inches(9), Inches(est_height)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = LIGHT_BG
            shape.line.color.rgb = RGBColor(200, 200, 200)
            
            tf = shape.text_frame
            tf.margin_left = Inches(0.2)
            tf.margin_top = Inches(0.1)
            
            fit_text_content(tf, intro_items, max_height_inches=est_height)
            top_cursor += est_height + 0.2
            
        # 3. Main Content Box (Bullets etc)
        if main_items:
            remaining_height = 7.0 - top_cursor
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(top_cursor), Inches(9), Inches(remaining_height))
            tf = content_box.text_frame
            tf.word_wrap = True
            fit_text_content(tf, main_items, max_height_inches=remaining_height)

        # Footer
        footer_box = slide.shapes.add_textbox(Inches(8), Inches(7.1), Inches(2), Inches(0.5))
        fp = footer_box.text_frame.add_paragraph()
        fp.text = "OTTOBITE 2026"
        fp.font.size = Pt(10)
        fp.font.color.rgb = RGBColor(150, 150, 150)
        fp.alignment = PP_ALIGN.RIGHT

    prs.save("OTTOBITE_Sunum_Final_Professional.pptx")
    print("Professional V8 Generated.")

if __name__ == "__main__":
    create_presentation()
