# encoding: utf-8
import json
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    BRICK_RED = RGBColor(192, 57, 43)
    DARK_GRAY = RGBColor(44, 62, 80)
    
    with open('structured_data_v3.json', 'r', encoding='utf-8') as f:
        slides_data = json.load(f)
    
    # Canvas Layout
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title Slide
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    p = txBox.text_frame.add_paragraph()
    p.text = "OTTOBITE"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = BRICK_RED
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1))
    p2 = txBox2.text_frame.add_paragraph()
    p2.text = "Garson Davranışları ve İş Önceliği Rehberi"
    p2.font.size = Pt(28)
    p2.font.color.rgb = DARK_GRAY
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = "Calibri Light"
    
    # Content Slides
    for slide_info in slides_data:
        if "OTTOBITE Garson Rehberi" in slide_info["title"] and len(slide_info["content"]) < 2:
            continue
            
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Header
        title_text = slide_info['title']
        # Fix numbers spacing "9 . " -> "9."
        title_text = re.sub(r'(\d)\s+\.', r'\1.', title_text)
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
        tp = title_box.text_frame.add_paragraph()
        tp.text = title_text
        tp.font.size = Pt(32)
        tp.font.bold = True
        tp.font.color.rgb = BRICK_RED
        tp.font.name = "Arial"
        
        # Separator
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(9), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = BRICK_RED
        line.line.fill.background()
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for item in slide_info['content']:
            text_val = item['text']
            item_type = item['type']
            
            p = tf.add_paragraph()
            
            # Styling based on type
            if item_type == "subheader":
                p.text = text_val
                p.font.bold = True
                p.space_before = Pt(12)
                p.font.color.rgb = DARK_GRAY
            
            elif item_type == "bullet":
                p.text = "• " + text_val
                p.level = 0 # Indent?
                # Handled by "• " manual bullet for better control on textbox
            
            else: # Normal text or emphasis
                if "kabul edilemez" in text_val.lower():
                    p.text = text_val
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(192, 0, 0) # Red emphasis
                else:
                    p.text = text_val
            
            p.font.size = Pt(18)
            p.font.name = "Calibri"
            if item_type != "subheader":
                p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(8)

        # Footer
        footer_box = slide.shapes.add_textbox(Inches(8), Inches(7.1), Inches(2), Inches(0.5))
        fp = footer_box.text_frame.add_paragraph()
        fp.text = "OTTOBITE 2026"
        fp.font.size = Pt(10)
        fp.font.color.rgb = RGBColor(150, 150, 150)
        fp.alignment = PP_ALIGN.RIGHT
    
    prs.save("OTTOBITE_Sunum_Final.pptx")
    print("Final Presentation generated.")

if __name__ == "__main__":
    create_presentation()
