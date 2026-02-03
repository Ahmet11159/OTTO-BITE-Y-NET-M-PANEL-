# encoding: utf-8
import json
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def final_heal(text):
    # Final healing for presentation display
    replacements = [
        (r'da\s*vr\s*anış', 'Davranış'),
        (r'mis\s*afir', 'Misafir'),
        (r's\s*er\s*vis', 'Servis'),
        (r'k\s*ur\s*al', 'Kural'),
        (r'bağ\s*lı', 'bağlı'),
        (r'g\s*er\s*ek\s*en', 'gereken'),
        (r'i\s*l\s*et\s*iş\s*im', 'İletişim'),
        (r'y\s*ak\s*l\s*aş\s*ma', 'Yaklaşma'),
        (r'y\s*üz\s*üne', 'yüzüne'),
        (r'k\s*ullanıl', 'Kullanıl'),
        (r's\s*an\s*iy\s*e', 'Saniye'),
        (r'ilgileniy\s*or\s*um', 'İlgileniyorum'),
        (r'yardımcı\s*oluy\s*or\s*um', 'Yardımcı Oluyorum'),
        (r'i\s*s\s*t\s*er', 'İster'),
        (r'd\s*ikk\s*at', 'Dikkat'),
        (r't\s*ak\s*ım', 'Takım'),
        (r'b\s*ilinci', 'Bilinci'),
        (r'i\s*s\s*t\s*er', 'İster'),
        (r'u\s*ygul\s*anmak', 'Uygulanmak'),
        (r'z\s*or\s*undadır', 'Zorundadır'),
        (r'h\s*izmet', 'Hizmet'),
        (r'g\s*ecik\s*me\s*si', 'Gecikmesi'),
        (r'\.\.', '.'),
    ]
    for pattern, replacement in replacements:
        text = re.sub(pattern, replacement, text, flags=re.IGNORECASE)
    return text

def create_presentation():
    prs = Presentation()
    
    # Colors
    BRICK_RED = RGBColor(192, 57, 43) # #C0392B
    DARK_GRAY = RGBColor(44, 62, 80)  # #2C3E50
    LIGHT_GRAY = RGBColor(236, 240, 241) # #ECF0F1 (Background hint if possible, usually plain white is safer)
    
    with open('structured_data_v2.json', 'r', encoding='utf-8') as f:
        slides_data = json.load(f)
    
    # ==========================
    # 1. TITLE SLIDE
    # ==========================
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    
    # Background (White default)
    
    # Title Text
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "OTTOBITE"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = BRICK_RED
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    
    # Subtitle
    top = Inches(3.8)
    txBox2 = slide.shapes.add_textbox(left, top, width, height)
    tf2 = txBox2.text_frame
    p2 = tf2.add_paragraph()
    p2.text = "Garson Davranışları ve İş Önceliği Rehberi"
    p2.font.size = Pt(28)
    p2.font.color.rgb = DARK_GRAY
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = "Calibri Light"
    
    # Decorative Line
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3), Inches(3.6), Inches(4), Inches(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRICK_RED
    shape.line.fill.background() # No line border
    
    # ==========================
    # 2. CONTENT SLIDES
    # ==========================
    # Custom Layout logic on Blank Slide for full control
    for slide_info in slides_data:
        # Skip if title looks like main title intro
        if "OTTOBITE Garson Rehberi" in slide_info["title"] and len(slide_info["content"]) < 2:
            continue
            
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
        
        # --- Header ---
        title_text = final_heal(slide_info['title'])
        # Fix repeated numbers "3 4."
        title_text = re.sub(r'(\d)\s+(\d)', r'\1\2', title_text)
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
        tp = title_box.text_frame.add_paragraph()
        tp.text = title_text
        tp.font.size = Pt(36)
        tp.font.bold = True
        tp.font.color.rgb = BRICK_RED
        tp.font.name = "Arial"
        
        # --- Separator Line ---
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(9), Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = BRICK_RED # Brick Red Line
        line.line.fill.background()
        
        # --- Content Body ---
        top_pos = 1.6
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(top_pos), Inches(9), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for item in slide_info['content']:
            cleaned_item = final_heal(item)
            cleaned_item = re.sub(r'(\d)\s+(\d)', r'\1\2', cleaned_item)
            
            p = tf.add_paragraph()
            p.text = "• " + cleaned_item # Custom Bullet
            p.font.size = Pt(20)
            p.font.name = "Calibri"
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(12) # Better spacing
            p.space_before = Pt(6)

        # --- Footer ---
        footer_box = slide.shapes.add_textbox(Inches(8), Inches(7.1), Inches(2), Inches(0.5))
        fp = footer_box.text_frame.add_paragraph()
        fp.text = "OTTOBITE 2026"
        fp.font.size = Pt(10)
        fp.font.color.rgb = RGBColor(150, 150, 150)
        fp.alignment = PP_ALIGN.RIGHT
    
    prs.save("OTTOBITE_Sunum_Modern.pptx")
    print("Modern Presentation generated: OTTOBITE_Sunum_Modern.pptx")

if __name__ == "__main__":
    create_presentation()
