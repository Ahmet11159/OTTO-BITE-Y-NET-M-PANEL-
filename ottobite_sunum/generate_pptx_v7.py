# encoding: utf-8
import json
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    BRICK_RED = RGBColor(192, 57, 43)
    DARK_GRAY = RGBColor(44, 62, 80)
    
    with open('structured_data_final.json', 'r', encoding='utf-8') as f:
        slides_data = json.load(f)
    
    # --- Helper: Add Title Slide ---
    def add_title_slide(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
        p = txBox.text_frame.add_paragraph()
        p.text = "OTTOBITE"
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = BRICK_RED
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Arial"
        
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1))
        p2 = txBox2.text_frame.add_paragraph()
        p2.text = "Garson Davranışları ve İş Önceliği Rehberi"
        p2.font.size = Pt(28)
        p2.font.color.rgb = DARK_GRAY
        p2.alignment = PP_ALIGN.CENTER
        p2.font.name = "Calibri Light"
        
        # Line
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(3.6), Inches(4), Inches(0.05))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BRICK_RED
        shape.line.fill.background()

    # --- Helper: Add Content Slide with Auto-Scaling ---
    def add_content_slide(prs, title, content_items):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Header
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
        tp = title_box.text_frame.add_paragraph()
        tp.text = title
        tp.font.size = Pt(32)
        tp.font.bold = True
        tp.font.color.rgb = BRICK_RED
        tp.font.name = "Arial"
        
        # Line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(9), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = BRICK_RED
        line.line.fill.background()
        
        # --- Font Size Calculation ---
        # Heuristic: Count items and total length
        total_chars = sum(len(i['text']) for i in content_items)
        num_items = len(content_items)
        
        # Base settings (Standard)
        font_size = 20
        subheader_size = 22
        space_after = 8
        space_before_subheader = 14
        
        # Scaling Logic
        if num_items > 18 or total_chars > 1100:
            font_size = 12
            subheader_size = 14
            space_after = 3
            space_before_subheader = 6
        elif num_items > 14 or total_chars > 850:
            font_size = 14
            subheader_size = 16
            space_after = 4
            space_before_subheader = 8
        elif num_items > 10 or total_chars > 600:
            font_size = 16
            subheader_size = 18
            space_after = 6
            space_before_subheader = 10
            
        # Content Box
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(5.2))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for item in content_items:
            text_val = item['text']
            item_type = item['type']
            
            p = tf.add_paragraph()
            
            if item_type == "subheader":
                p.text = text_val
                p.font.bold = True
                p.font.size = Pt(subheader_size)
                p.space_before = Pt(space_before_subheader)
                p.font.color.rgb = DARK_GRAY
            elif item_type == "bullet":
                p.text = "• " + text_val
                p.font.size = Pt(font_size)
                p.space_after = Pt(space_after)
            elif item_type == "emphasis":
                p.text = text_val
                p.font.bold = True
                p.font.color.rgb = RGBColor(192, 0, 0)
                p.font.size = Pt(font_size + 1)
                p.space_before = Pt(space_after)
            else:
                p.text = text_val
                p.font.size = Pt(font_size)
                p.space_after = Pt(space_after)
            
            p.font.name = "Calibri"
            if item_type != "emphasis" and item_type != "subheader":
                p.font.color.rgb = DARK_GRAY

        # Footer
        footer_box = slide.shapes.add_textbox(Inches(8), Inches(7.1), Inches(2), Inches(0.5))
        fp = footer_box.text_frame.add_paragraph()
        fp.text = "OTTOBITE 2026"
        fp.font.size = Pt(10)
        fp.font.color.rgb = RGBColor(150, 150, 150)
        fp.alignment = PP_ALIGN.RIGHT

    add_title_slide(prs)
    
    for slide_info in slides_data:
        if len(slide_info['content']) == 0: continue
        if slide_info['title'] == "OTTOBITE Garson Rehberi" and len(slide_info['content']) < 3: continue
        
        add_content_slide(prs, slide_info['title'], slide_info['content'])

    prs.save("OTTOBITE_Sunum_Final_V7.pptx")
    print("V7 Presentation generated (Scaled Fonts).")

if __name__ == "__main__":
    create_presentation()
