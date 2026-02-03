# encoding: utf-8
"""
OTTOBITE V5 Enhanced - Preserves original V5 design with improvements:
1. Dynamic separator lines (proportional to title)
2. Comment windows for intro_box items
3. Content visibility (font scaling + smart splitting)
4. Full professional polish
"""

import json
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    BRICK_RED = RGBColor(192, 57, 43)
    DARK_GRAY = RGBColor(44, 62, 80)
    LIGHT_BG = RGBColor(245, 247, 250)
    BORDER_COLOR = RGBColor(215, 215, 215)
    
    with open('structured_data_final_v7.json', 'r', encoding='utf-8') as f:
        slides_data = json.load(f)
    
    # ==========================================
    # TITLE SLIDE (Preserved from V5)
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    p = txBox.text_frame.add_paragraph()
    p.text = "OTTOBITE"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = BRICK_RED
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1))
    p2 = txBox2.text_frame.add_paragraph()
    p2.text = "Garson Davranışları ve İş Önceliği Rehberi"
    p2.font.size = Pt(28)
    p2.font.color.rgb = DARK_GRAY
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = "Calibri Light"
    
    # Decorative Line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(3.6), Inches(4), Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRICK_RED
    shape.line.fill.background()
    
    # ==========================================
    # HELPER: Estimate content height
    # ==========================================
    def estimate_content_height(content_items, font_size):
        """Estimate total height in inches"""
        total_points = 0
        for item in content_items:
            text_len = len(item['text'])
            chars_per_line = int(650 / (0.5 * font_size))
            lines = max(1, math.ceil(text_len / chars_per_line))
            line_height = font_size * 1.3
            spacing = font_size * 0.4
            
            if item['type'] == 'subheader':
                spacing += font_size * 0.8
            
            total_points += (lines * line_height) + spacing
        
        return total_points / 72.0
    
    # ==========================================
    # HELPER: Add content slide
    # ==========================================
    def add_content_slide(prs, title, content_items, is_continuation=False):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        display_title = f"{title} (Devam)" if is_continuation else title
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
        tp = title_box.text_frame.add_paragraph()
        tp.text = display_title
        tp.font.size = Pt(32)
        tp.font.bold = True
        tp.font.color.rgb = BRICK_RED
        tp.font.name = "Arial"
        
        # --- DYNAMIC SEPARATOR LINE (New Feature) ---
        char_count = len(display_title)
        line_width = min(9.0, max(2.0, char_count * 0.20))
        
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0.5), Inches(1.3), 
            Inches(line_width), Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = BRICK_RED
        line.line.fill.background()
        
        # --- Separate intro_box from main content ---
        intro_items = [i for i in content_items if i['type'] == 'intro_box']
        main_items = [i for i in content_items if i['type'] != 'intro_box']
        
        current_top = 1.6
        
        # --- COMMENT WINDOW for intro items (New Feature) ---
        if intro_items:
            total_chars = sum(len(i['text']) for i in intro_items)
            num_items = len(intro_items)
            est_lines = max(num_items, math.ceil(total_chars / 85))
            box_height = min(1.5, max(0.55, est_lines * 0.28 + 0.2))
            
            comment_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(current_top),
                Inches(9), Inches(box_height)
            )
            comment_box.fill.solid()
            comment_box.fill.fore_color.rgb = LIGHT_BG
            comment_box.line.color.rgb = BORDER_COLOR
            comment_box.line.width = Pt(1)
            
            tf = comment_box.text_frame
            tf.margin_left = Inches(0.2)
            tf.margin_right = Inches(0.2)
            tf.margin_top = Inches(0.1)
            tf.margin_bottom = Inches(0.1)
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP
            
            for idx, item in enumerate(intro_items):
                p = tf.add_paragraph()
                p.text = item['text']
                p.font.size = Pt(16)
                p.font.italic = True
                p.font.color.rgb = RGBColor(60, 60, 60)
                p.font.name = "Calibri"
                if idx > 0:
                    p.space_before = Pt(4)
            
            current_top += box_height + 0.2
        
        # --- MAIN CONTENT (Preserved V5 style with auto-fit) ---
        if main_items:
            available_height = 6.9 - current_top
            
            # Calculate appropriate font size
            sizes = [18, 16, 15, 14, 13, 12]
            selected_size = 12
            
            for size in sizes:
                est_height = estimate_content_height(main_items, size)
                if est_height <= available_height:
                    selected_size = size
                    break
            
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(current_top),
                Inches(9), Inches(available_height)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP
            
            for item in main_items:
                text_val = item['text']
                item_type = item['type']
                
                p = tf.add_paragraph()
                p.font.name = "Calibri"
                
                if item_type == "subheader":
                    p.text = text_val
                    p.font.bold = True
                    p.font.size = Pt(min(selected_size + 2, 20))
                    p.space_before = Pt(12)
                    p.space_after = Pt(4)
                    p.font.color.rgb = DARK_GRAY
                    
                elif item_type == "bullet":
                    p.text = "• " + text_val
                    p.font.size = Pt(selected_size)
                    p.space_after = Pt(4)
                    p.font.color.rgb = DARK_GRAY
                    
                elif item_type == "emphasis":
                    p.text = text_val
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(192, 0, 0)
                    p.font.size = Pt(selected_size + 1)
                    p.space_before = Pt(6)
                    p.space_after = Pt(4)
                    
                else:  # body_text
                    p.text = text_val
                    p.font.size = Pt(selected_size)
                    p.space_after = Pt(6)
                    p.font.color.rgb = DARK_GRAY
        
        # Footer (Preserved from V5)
        footer_box = slide.shapes.add_textbox(Inches(8), Inches(7.1), Inches(2), Inches(0.5))
        fp = footer_box.text_frame.add_paragraph()
        fp.text = "OTTOBITE 2026"
        fp.font.size = Pt(10)
        fp.font.color.rgb = RGBColor(150, 150, 150)
        fp.alignment = PP_ALIGN.RIGHT
    
    # ==========================================
    # PROCESS SLIDES
    # ==========================================
    MAX_CONTENT_LINES = 16
    
    for slide_info in slides_data:
        if len(slide_info['content']) == 0:
            continue
        if slide_info['title'] == "OTTOBITE Garson Rehberi" and len(slide_info['content']) < 3:
            continue
        
        title = slide_info['title']
        all_content = slide_info['content']
        
        # Estimate if content fits
        def count_lines(items):
            total = 0
            for item in items:
                text_len = len(item['text'])
                lines = max(1, math.ceil(text_len / 75))
                if item['type'] == 'subheader':
                    lines += 0.5
                total += lines
            return total
        
        total_lines = count_lines(all_content)
        
        if total_lines <= MAX_CONTENT_LINES:
            add_content_slide(prs, title, all_content, is_continuation=False)
        else:
            # Split content
            chunks = []
            current_chunk = []
            current_lines = 0
            
            for item in all_content:
                text_len = len(item['text'])
                item_lines = max(1, math.ceil(text_len / 75))
                if item['type'] == 'subheader':
                    item_lines += 0.5
                
                if current_lines + item_lines > MAX_CONTENT_LINES and current_chunk:
                    chunks.append(current_chunk)
                    current_chunk = []
                    current_lines = 0
                
                current_chunk.append(item)
                current_lines += item_lines
            
            if current_chunk:
                chunks.append(current_chunk)
            
            for idx, chunk in enumerate(chunks):
                add_content_slide(prs, title, chunk, is_continuation=(idx > 0))
    
    # Save
    prs.save("OTTOBITE_Sunum_V5_Enhanced.pptx")
    print("✓ V5 Enhanced Presentation generated successfully!")

if __name__ == "__main__":
    create_presentation()
