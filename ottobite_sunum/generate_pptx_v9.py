# encoding: utf-8
import json
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR

def create_presentation():
    prs = Presentation()
    BRICK_RED = RGBColor(192, 57, 43)
    DARK_GRAY = RGBColor(44, 62, 80)
    LIGHT_BG = RGBColor(240, 242, 245) # Soft Gray for comments
    BORDER_COLOR = RGBColor(200, 200, 200)
    
    with open('structured_data_final_v7.json', 'r', encoding='utf-8') as f:
        slides_data = json.load(f)

    # --- Auto-Fit Helper ---
    def fit_text_content(tf, content_items, max_height_inches=5.8):
        sizes = [20, 18, 16, 14, 12, 11]
        selected_size = 12
        
        for size in sizes:
            current_height = 0
            line_height = size * 1.3
            para_spacing = size * 0.5
            total_points = 0
            
            for item in content_items:
                text = item['text']
                max_chars = 900 / (0.55 * size)
                lines = math.ceil(len(text) / max_chars)
                if lines < 1: lines = 1
                
                item_h = (lines * line_height) + para_spacing
                if item['type'] == 'subheader': item_h += 10
                
                total_points += item_h
            
            total_inches = total_points / 72
            if total_inches <= max_height_inches:
                selected_size = size
                break
        
        # Render
        for item in content_items:
            text_val = item['text']
            item_type = item['type']
            
            p = tf.add_paragraph()
            
            if item_type == "subheader":
                p.text = text_val
                p.font.bold = True
                p.font.size = Pt(selected_size + 2)
                p.font.color.rgb = DARK_GRAY
                p.space_before = Pt(selected_size)
            elif item_type == "bullet":
                p.text = "• " + text_val
                p.font.size = Pt(selected_size)
                p.font.color.rgb = DARK_GRAY
                p.space_after = Pt(selected_size * 0.4)
            elif item_type == "emphasis":
                p.text = text_val.upper()
                p.font.bold = True
                p.font.color.rgb = RGBColor(192, 0, 0)
                p.font.size = Pt(selected_size)
                p.space_before = Pt(selected_size * 0.5)
                p.alignment = PP_ALIGN.CENTER
            else:
                 p.text = text_val
                 p.font.size = Pt(selected_size)
                 p.font.color.rgb = DARK_GRAY
                 p.space_after = Pt(selected_size * 0.4)
            
            p.font.name = "Calibri"

    # --- Title Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    p = txBox.text_frame.add_paragraph()
    p.text = "OTTOBITE"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = BRICK_RED
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1))
    p2 = txBox2.text_frame.add_paragraph()
    p2.text = "Garson Davranışları ve İş Önceliği Rehberi"
    p2.font.size = Pt(28)
    p2.font.color.rgb = DARK_GRAY
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = "Calibri Light"
    
    # --- Content Slides ---
    for slide_info in slides_data:
        if len(slide_info['content']) == 0: continue
        if slide_info['title'] == "OTTOBITE Garson Rehberi" and len(slide_info['content']) < 3: continue
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 1. Header
        title_text = slide_info['title']
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
        tp = title_box.text_frame.add_paragraph()
        tp.text = title_text
        tp.font.size = Pt(32)
        tp.font.bold = True
        tp.font.color.rgb = BRICK_RED
        tp.font.name = "Arial"
        
        # Dynamic Line
        char_count = len(title_text)
        est_width = char_count * 0.23
        if est_width > 9: est_width = 9
        if est_width < 1: est_width = 1
        
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), est_width, Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = BRICK_RED
        line.line.fill.background()
        
        # 2. Separation Logic: Intro Box vs Main Content
        intro_items = []
        main_items = []
        
        for item in slide_info['content']:
            if item['type'] == 'intro_box':
                intro_items.append(item)
            else:
                main_items.append(item)
        
        top_cursor = 1.6
        
        # 3. Render Intro Box (if any)
        if intro_items:
            # Calc Height
            est_lines = sum([math.ceil(len(i['text'])/90) for i in intro_items])
            est_height = (est_lines * 0.35) + 0.4 # padding
            if est_height < 0.8: est_height = 0.8
            
            # Draw Shape
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, 
                Inches(0.5), Inches(top_cursor), Inches(9), Inches(est_height)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = LIGHT_BG
            shape.line.color.rgb = BORDER_COLOR
            shape.line.width = Pt(1)
            
            tf = shape.text_frame
            tf.margin_left = Inches(0.2)
            tf.margin_right = Inches(0.2)
            tf.margin_top = Inches(0.15)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            for item in intro_items:
                p = tf.add_paragraph()
                p.text = item['text']
                p.font.size = Pt(16)
                p.font.italic = True
                p.font.color.rgb = RGBColor(80, 80, 80)
                p.space_after = Pt(6)
                p.alignment = PP_ALIGN.LEFT
            
            top_cursor += est_height + 0.2 # Space below box
            
        # 4. Render Main Content (Bullets etc)
        if main_items:
            remaining_height = 7.0 - top_cursor
            if remaining_height < 1: remaining_height = 1
            
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(top_cursor), Inches(9), Inches(remaining_height))
            tf = content_box.text_frame
            tf.word_wrap = True
            fit_text_content(tf, main_items, max_height_inches=remaining_height)

        # Footer
        footer_box = slide.shapes.add_textbox(Inches(8), Inches(7.1), Inches(2), Inches(0.5))
        fp = footer_box.text_frame.add_paragraph()
        fp.text = "OTTOBITE 2026"
        fp.font.size = Pt(10)
        fp.font.color.rgb = RGBColor(150, 150, 150)
        fp.alignment = PP_ALIGN.RIGHT

    prs.save("OTTOBITE_Sunum_Final_V9_Windows.pptx")
    print("V9 Presentation generated.")

if __name__ == "__main__":
    create_presentation()
