# encoding: utf-8
import json
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    BRICK_RED = RGBColor(192, 57, 43)
    DARK_GRAY = RGBColor(44, 62, 80)
    
    with open('structured_data_final.json', 'r', encoding='utf-8') as f:
        slides_data = json.load(f)
    
    # --- Helper: Add Title Slide ---
    def add_title_slide(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
        p = txBox.text_frame.add_paragraph()
        p.text = "OTTOBITE"
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = BRICK_RED
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Arial"
        
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1))
        p2 = txBox2.text_frame.add_paragraph()
        p2.text = "Garson Davranışları ve İş Önceliği Rehberi"
        p2.font.size = Pt(28)
        p2.font.color.rgb = DARK_GRAY
        p2.alignment = PP_ALIGN.CENTER
        p2.font.name = "Calibri Light"
        
        # Line
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(3.6), Inches(4), Inches(0.05))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BRICK_RED
        shape.line.fill.background()

    # --- Helper: Add Content Slide ---
    def add_content_slide(prs, title, content_items, is_continued=False):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Header
        disp_title = title + " (Devam)" if is_continued else title
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
        tp = title_box.text_frame.add_paragraph()
        tp.text = disp_title
        tp.font.size = Pt(32)
        tp.font.bold = True
        tp.font.color.rgb = BRICK_RED
        tp.font.name = "Arial"
        
        # Line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(9), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = BRICK_RED
        line.line.fill.background()
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for item in content_items:
            text_val = item['text']
            item_type = item['type']
            
            p = tf.add_paragraph()
            
            if item_type == "subheader":
                p.text = text_val
                p.font.bold = True
                p.font.size = Pt(20)
                p.space_before = Pt(14)
                p.font.color.rgb = DARK_GRAY
            elif item_type == "bullet":
                p.text = "• " + text_val
                p.font.size = Pt(18)
                p.space_after = Pt(6)
            elif item_type == "emphasis":
                p.text = text_val
                p.font.bold = True
                p.font.color.rgb = RGBColor(192, 0, 0)
                p.font.size = Pt(19)
                p.space_before = Pt(8)
            else:
                p.text = text_val
                p.font.size = Pt(18)
                p.space_after = Pt(8)
            
            p.font.name = "Calibri"
            if item_type != "emphasis" and item_type != "subheader":
                p.font.color.rgb = DARK_GRAY

        # Footer
        footer_box = slide.shapes.add_textbox(Inches(8), Inches(7.1), Inches(2), Inches(0.5))
        fp = footer_box.text_frame.add_paragraph()
        fp.text = "OTTOBITE 2026"
        fp.font.size = Pt(10)
        fp.font.color.rgb = RGBColor(150, 150, 150)
        fp.alignment = PP_ALIGN.RIGHT

    # --- Main Logic ---
    add_title_slide(prs)
    
    MAX_LINES = 12 # Approximate conservative limit
    
    for slide_info in slides_data:
        if len(slide_info['content']) == 0: continue
        if slide_info['title'] == "OTTOBITE Garson Rehberi" and len(slide_info['content']) < 3: continue
        
        title = slide_info['title']
        content = slide_info['content']
        
        # Calculate chunks
        current_chunk = []
        current_cost = 0
        
        for item in content:
            # Estimate cost
            cost = 1
            if len(item['text']) > 80: cost = 2
            if item['type'] == 'subheader': cost = 2
            
            if current_cost + cost > MAX_LINES:
                # Flush
                add_content_slide(prs, title, current_chunk, is_continued=(len(current_chunk) != len(content) and len(current_chunk) > 0))
                # If this is the second or third chunk, the title logic handles "Devam" if we passed is_continued=True? 
                # Wait, "is_continued" param needs to be stateful.
                
                # Correction: "is_continued" should be True if we have already flushed nicely.
                # Actually, simpler: Flush current, start new. 
                # First chunk: Not continued.
                # Subsequent chunks: Continued.
                title = slide_info['title'] # Reset title base, but we need to mark it.
                
                # Let's fix loop logic.
                pass 
                
        # Re-do Loop Logic for splitting
        chunks = []
        chunk = []
        cost = 0
        
        for item in content:
            item_cost = 1
            if len(item['text']) > 80: item_cost += 1
            if item['type'] == 'subheader': item_cost += 1
            
            if cost + item_cost > MAX_LINES:
                chunks.append(chunk)
                chunk = []
                cost = 0
            
            chunk.append(item)
            cost += item_cost
            
        if chunk:
            chunks.append(chunk)
            
        # Create slides from chunks
        for i, chk in enumerate(chunks):
            is_cont = (i > 0)
            add_content_slide(prs, title, chk, is_continued=is_cont)

    prs.save("OTTOBITE_Sunum_Final_V6.pptx")
    print("V6 Presentation generated.")

if __name__ == "__main__":
    create_presentation()
