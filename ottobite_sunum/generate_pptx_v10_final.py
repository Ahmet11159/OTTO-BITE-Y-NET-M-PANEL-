# encoding: utf-8
import json
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR

def create_presentation():
    prs = Presentation()
    
    # Professional Color Palette
    BRICK_RED = RGBColor(192, 57, 43)      # Primary accent
    DARK_GRAY = RGBColor(44, 62, 80)       # Body text
    LIGHT_BG = RGBColor(245, 247, 250)     # Comment window background
    BORDER_COLOR = RGBColor(220, 220, 220) # Subtle borders
    FOOTER_GRAY = RGBColor(150, 150, 150)  # Footer text
    
    with open('structured_data_final_v7.json', 'r', encoding='utf-8') as f:
        slides_data = json.load(f)

    # --- Enhanced Auto-Fit Algorithm ---
    def fit_text_content(tf, content_items, max_height_inches=5.5):
        """
        Intelligently scales font size to fit content within available space.
        Uses character count and line wrapping estimation.
        """
        sizes = [20, 18, 16, 15, 14, 13, 12, 11]
        selected_size = 11  # Safe fallback
        
        for size in sizes:
            line_height = size * 1.25
            base_spacing = max(3, size * 0.3)
            total_points = 0
            
            for item in content_items:
                text = item['text']
                # Estimate characters per line (9 inches ≈ 648 points)
                chars_per_line = int(900 / (0.52 * size))
                
                # Calculate wrapped lines
                lines_needed = max(1, math.ceil(len(text) / chars_per_line))
                
                # Calculate item height with type-specific spacing
                item_height = lines_needed * line_height
                
                if item['type'] == 'subheader':
                    item_height += size * 1.2  # Extra space before
                elif item['type'] == 'emphasis':
                    item_height += size * 0.8
                
                item_height += base_spacing
                total_points += item_height
            
            # Convert to inches and check if it fits
            total_inches = total_points / 72.0
            if total_inches <= max_height_inches:
                selected_size = size
                break
        
        # Render content with selected size
        for item in content_items:
            text_val = item['text']
            item_type = item['type']
            
            p = tf.add_paragraph()
            p.font.name = "Calibri"
            
            if item_type == "subheader":
                p.text = text_val
                p.font.bold = True
                p.font.size = Pt(min(selected_size + 3, 22))
                p.font.color.rgb = DARK_GRAY
                p.space_before = Pt(selected_size * 1.2)
                p.space_after = Pt(selected_size * 0.3)
                
            elif item_type == "bullet":
                p.text = "• " + text_val
                p.font.size = Pt(selected_size)
                p.font.color.rgb = DARK_GRAY
                p.space_after = Pt(max(3, selected_size * 0.3))
                
            elif item_type == "emphasis":
                p.text = text_val.upper()
                p.font.bold = True
                p.font.color.rgb = RGBColor(192, 0, 0)
                p.font.size = Pt(selected_size + 1)
                p.space_before = Pt(selected_size * 0.6)
                p.space_after = Pt(selected_size * 0.4)
                p.alignment = PP_ALIGN.CENTER
                
            else:  # body_text
                p.text = text_val
                p.font.size = Pt(selected_size)
                p.font.color.rgb = DARK_GRAY
                p.space_after = Pt(max(4, selected_size * 0.4))

    # ==========================================
    # TITLE SLIDE
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Main title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    p = txBox.text_frame.add_paragraph()
    p.text = "OTTOBITE"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = BRICK_RED
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    
    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1))
    p2 = txBox2.text_frame.add_paragraph()
    p2.text = "Garson Davranışları ve İş Önceliği Rehberi"
    p2.font.size = Pt(28)
    p2.font.color.rgb = DARK_GRAY
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = "Calibri Light"
    
    # Decorative accent line
    accent = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.5), Inches(3.55), Inches(3), Inches(0.08)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = BRICK_RED
    accent.line.fill.background()
    
    # ==========================================
    # CONTENT SLIDES
    # ==========================================
    for slide_info in slides_data:
        # Skip empty or intro-only slides
        if len(slide_info['content']) == 0:
            continue
        if slide_info['title'] == "OTTOBITE Garson Rehberi" and len(slide_info['content']) < 3:
            continue
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # --- HEADER with DYNAMIC SEPARATOR LINE ---
        title_text = slide_info['title']
        
        # Title text box
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.9))
        tp = title_box.text_frame.add_paragraph()
        tp.text = title_text
        tp.font.size = Pt(32)
        tp.font.bold = True
        tp.font.color.rgb = BRICK_RED
        tp.font.name = "Arial"
        tp.alignment = PP_ALIGN.LEFT
        
        # Dynamic separator line (proportional to title length)
        title_length = len(title_text)
        # Estimate: ~0.22 inches per character for Arial Bold 32pt
        line_width = min(9.0, max(2.0, title_length * 0.22))
        
        separator = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.28), Inches(line_width), Pt(3)
        )
        separator.fill.solid()
        separator.fill.fore_color.rgb = BRICK_RED
        separator.line.fill.background()
        
        # --- CONTENT SEPARATION ---
        intro_items = []
        main_items = []
        
        for item in slide_info['content']:
            if item['type'] == 'intro_box':
                intro_items.append(item)
            else:
                main_items.append(item)
        
        vertical_cursor = 1.55
        
        # --- COMMENT WINDOW (if intro items exist) ---
        if intro_items:
            # Estimate height needed
            total_chars = sum(len(item['text']) for item in intro_items)
            num_items = len(intro_items)
            estimated_lines = max(num_items, math.ceil(total_chars / 85))
            box_height = min(1.8, max(0.7, estimated_lines * 0.3 + 0.35))
            
            # Create rounded rectangle
            comment_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(vertical_cursor), Inches(9), Inches(box_height)
            )
            comment_box.fill.solid()
            comment_box.fill.fore_color.rgb = LIGHT_BG
            comment_box.line.color.rgb = BORDER_COLOR
            comment_box.line.width = Pt(1.5)
            
            # Add text to comment box
            tf = comment_box.text_frame
            tf.margin_left = Inches(0.25)
            tf.margin_right = Inches(0.25)
            tf.margin_top = Inches(0.18)
            tf.margin_bottom = Inches(0.18)
            tf.vertical_anchor = MSO_ANCHOR.TOP
            tf.word_wrap = True
            
            font_size = 15 if box_height > 1.2 else 14
            
            for idx, item in enumerate(intro_items):
                p = tf.add_paragraph()
                p.text = item['text']
                p.font.size = Pt(font_size)
                p.font.italic = True
                p.font.color.rgb = RGBColor(70, 70, 70)
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.LEFT
                if idx < len(intro_items) - 1:
                    p.space_after = Pt(8)
            
            vertical_cursor += box_height + 0.25
        
        # --- MAIN CONTENT AREA ---
        if main_items:
            available_height = 7.0 - vertical_cursor
            available_height = max(1.0, available_height)  # Ensure minimum space
            
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(vertical_cursor),
                Inches(9), Inches(available_height)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP
            
            fit_text_content(tf, main_items, max_height_inches=available_height - 0.1)
        
        # --- FOOTER ---
        footer_box = slide.shapes.add_textbox(Inches(8.2), Inches(7.15), Inches(1.3), Inches(0.35))
        fp = footer_box.text_frame.add_paragraph()
        fp.text = "OTTOBITE 2026"
        fp.font.size = Pt(10)
        fp.font.color.rgb = FOOTER_GRAY
        fp.alignment = PP_ALIGN.RIGHT
        fp.font.name = "Calibri"
    
    # ==========================================
    # SAVE PRESENTATION
    # ==========================================
    prs.save("OTTOBITE_Sunum_FINAL.pptx")
    print("✓ Ultra-Professional V10 Presentation Generated Successfully.")

if __name__ == "__main__":
    create_presentation()
