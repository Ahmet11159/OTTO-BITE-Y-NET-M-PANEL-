# encoding: utf-8
import json
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def final_heal(text):
    # JSON çıktısında görülen son hataların manuel düzeltilmesi
    replacements = [
        (r'k\s*ab\s*ul', 'kabul'),
        (r'edileme\s*z', 'edilemez'),
        (r'm\s*u\s*t\s*f\s*a\s*k\s*t\s*a\s*n', 'mutfaktan'),
        (r'k\s*a\s*s\s*a\s*y\s*a', 'kasaya'),
        (r's\s*er\s*gile\s*y\s*en', 'sergileyen'),
        (r'k\s*e\s*y\s*i\s*ﬂ\s*i', 'keyifli'),
        (r's\s*a\s*ğ\s*l\s*adık', 'sağladık'),
        (r'g\s*ör\s*üş\s*er\s*ek', 'görüşerek'),
        (r'b\s*a\s*ş\s*ın\s*a', 'başına'),
        (r'g\s*iz\s*lemek', 'gizlemek'),
        (r'i\s*s\s*tisna\s*i', 'istisnai'),
        (r'v\s*ur\s*g\s*ul\s*anar\s*ak', 'vurgulanarak'),
        (r'a\s*k\s*t\s*arılır', 'aktarılır'),
        (r'i\s*htiy\s*aç', 'ihtiyaç'),
        (r'ö\s*ncelik\s*l\s*idir', 'önceliklidir'),
        (r'o\s*na\s*yı', 'onayı'),
        (r'm\s*u\s*t\s*ab\s*ak\s*a\s*t', 'mutabakat'),
        (r'z\s*or\s*unludur', 'zorunludur'),
        (r'b\s*o\s*ş\s*l\s*u\s*k', 'boşluk'),
        (r'y\s*ar\s*a\s*t\s*m\s*a\s*y\s*ac\s*ak', 'yaratmayacak'),
        (r'b\s*a\s*ğ\s*l\s*antılı', 'bağlantılı'),
        (r'ç\s*ı\s*k\s*ama\s*z', 'çıkamaz'),
        (r'i\s*ler\s*leme\s*si', 'ilerlemesi'),
        (r'k\s*e\s*sin\s*tisiz', 'kesintisiz'),
        (r'G\s*ör\s*üş\s*ü', 'Görüşü'),
        (r'e\s*ng\s*el\s*le\s*y\s*ecek', 'engelleyecek'),
        (r'b\s*ar\s*dak\s*l\s*ar', 'bardaklar'),
        (r't\s*ab\s*ak\s*l\s*ar', 'tabaklar'),
        (r'k\s*ontr\s*ols\s*üz', 'kontrolsüz'),
        (r'y\s*ı\s*ğ\s*ıl\s*ama\s*z', 'yığılamaz'),
        (r'y\s*ı\s*ğ\s*m\s*a', 'yığma'),
        (r'k\s*a\s*y\s*mas\s*ına', 'kaymasına'),
        (r'd\s*i\s*z\s*i\s*m', 'dizim'),
        (r'a\s*ğ\s*ır\s*lık', 'ağırlık'),
        (r'm\s*er\s*k\s*e\s*z', 'merkez'),
        (r't\s*op\s*l\s*anmalıdır', 'toplanmalıdır'),
        (r'b\s*üy\s*ük', 'büyük'),
        (r'k\s*ü\s*ç\s*ü\s*ğ\s*e', 'küçüğe'),
        (r'i\s*s\s*tiﬂeme', 'istifleme'),
        (r'y\s*ap\s*ılmalıdır', 'yapılmalıdır'),
        (r'k\s*ır\s*ılma', 'kırılma'),
        (r'p\s*ar\s*ç\s*al\s*ar', 'parçalar'),
        (r'a\s*y\s*rış\s*tırılmalıdır', 'ayrıştırılmalıdır'),
        (r'y\s*er\s*le\s*ş\s*tirilir', 'yerleştirilir'),
        (r'k\s*apıdan', 'kapıda'),
        (r'g\s*ir\s*en', 'giren'),
        (r'f\s*ark', 'fark'),
        (r'e\s*dilmelidir', 'edilmelidir'),
        (r't\s*ebes\s*s\s*üm', 'tebessüm'),
        (r'k\s*ar\s*şıl\s*aması', 'karşılaması'),
        (r'y\s*apılır', 'yapılır'),
        (r'y\s*er\s*le\s*ş\s*ene', 'yerleşene'),
        (r'e\s*ş\s*lik', 'eşlik'),
        (r'e\s*dilir', 'edilir'),
        (r'b\s*ek\s*letileme\s*z', 'bekletilemez'),
        (r'g\s*el\s*ineme\s*z', 'gelinemez'),
        (r'i\s*z\s*ler', 'izler'),
        (r'i\s*ç\s*ec\s*ek\s*ler', 'içecekler'),
        (r'a\s*z\s*alan', 'azalan'),
        (r'p\s*e\s*ç\s*et\s*eler', 'peçeteler'),
        (r'b\s*o\s*ş\s*alan', 'boşalan'),
        (r'ç\s*e\s*vr\s*e\s*y\s*e', 'çevreye'),
        (r'b\s*a\s*k\s*ması', 'bakması'),
        (r'k\s*al\s*dırması', 'kaldırması'),
        (r'y\s*en\s*ilenme\s*si', 'yenilenmesi'),
        (r't\s*ek\s*lif', 'teklif'),
        (r'd\s*u\s*y\s*ulmadan', 'duyulmadan'),
        (r'd\s*es\s*t\s*e\s*ği', 'desteği'),
        (r's\s*a\s*ğ\s*l\s*anır', 'sağlanır'),
        (r's\s*ür\s*ek\s*li', 'sürekli'),
        (r't\s*ar\s*anır', 'taranır'),
        (r's\s*es\s*lenme\s*y\s*e', 'seslenmeye'),
        (r'b\s*ır\s*akmak', 'bırakmak'),
        (r's\s*ilme', 'silme'),
        (r'ç\s*ek\s*me', 'çekme'),
        (r'd\s*üz\s*enleme', 'düzenleme'),
        (r't\s*r\s*afiğine', 'trafiğine'),
        (r'b\s*ır\s*akıl\s*ar\s*ak', 'bırakılarak'),
        (r'o\s*dak\s*l\s*anıl\s*ama\s*z', 'odaklanılamaz'),
        (r'ç\s*aplı', 'çaplı'),
        (r'g\s*irilme\s*si', 'girilmesi'),
        (r's\s*at\s*ler\s*de', 'saatlerde'),
        (r't\s*am\s*aml\s*anır', 'tamamlanır'),
        (r'i\s*ş\s*lemleri', 'işlemleri'),
        (r'r\s*ah\s*atsız', 'rahatsız'),
        (r'e\s*tme\s*y\s*ecek', 'etmeyecek'),
        (r's\s*ess\s*iz\s*ce', 'sessizce'),
        (r'y\s*ür\s*ütülür', 'yürütülür'),
        (r'k\s*riz\s*e', 'krize'),
        (r'önceden', 'önceden'),
        (r'g\s*iderilir', 'giderilir'),
        # Add simpler ones
        (r'G\s*ars\s*on', 'Garson'),
        (r's\s*er\s*vis', 'servis'),
        (r'S\s*er\s*vis', 'Servis'),
        (r'O\s*TT\s*OBITE', 'OTTOBITE'),
    ]
    
    for pattern, replacement in replacements:
        text = re.sub(pattern, replacement, text, flags=re.IGNORECASE)
    
    return text

def create_presentation():
    prs = Presentation()
    
    # Kiremit Rengi (Brick Red)
    BRICK_RED = RGBColor(192, 57, 43) # #C0392B
    DARK_GRAY = RGBColor(44, 62, 80)  # #2C3E50
    
    # Load Structured Data
    with open('structured_data.json', 'r', encoding='utf-8') as f:
        slides_data = json.load(f)
    
    # 1. Title Slide
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "OTTOBITE"
    subtitle.text = "Garson Davranışları ve İş Önceliği Rehberi"
    
    # Style Title
    title.text_frame.paragraphs[0].font.color.rgb = BRICK_RED
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.name = "Arial"
    
    # 2. Content Slides
    blank_layout = prs.slide_layouts[1] # Title and Content
    
    for slide_info in slides_data:
        # Skip if title matches the main title (intro duplicate)
        if "OTTOBITE Garson Rehberi" in slide_info["title"]:
            # Maybe Intro slide content?
            # Let's check contents.
            pass
            
        slide = prs.slides.add_slide(blank_layout)
        
        # Title
        title_shape = slide.shapes.title
        raw_title = slide_info['title']
        clean_title = final_heal(raw_title)
        
        # Remove repeated numbering if messy (e.g. "3 4. Ort ak")
        # "3 4." -> "34."
        clean_title = re.sub(r'(\d)\s+(\d)', r'\1\2', clean_title)
        
        title_shape.text = clean_title
        title_shape.text_frame.paragraphs[0].font.color.rgb = BRICK_RED
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.name = "Arial"
        
        # Content
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear() # Clear default empty paragraph
        
        for item in slide_info['content']:
            p = tf.add_paragraph()
            cleaned_item = final_heal(item)
            cleaned_item = re.sub(r'(\d)\s+(\d)', r'\1\2', cleaned_item) # Fix numbers in content text too
            
            p.text = cleaned_item
            p.font.size = Pt(20)
            p.font.name = "Calibri"
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(10)
            
    # Save
    prs.save("OTTOBITE_Sunum.pptx")
    print("Presentation generated: OTTOBITE_Sunum.pptx")

if __name__ == "__main__":
    create_presentation()
