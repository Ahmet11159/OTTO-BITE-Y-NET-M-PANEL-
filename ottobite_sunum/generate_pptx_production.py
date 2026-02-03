# encoding: utf-8
"""
OTTOBITE Presentation Generator - Production Ready V11
Features:
- Content overflow: Splits into multiple slides (not shrink to unreadable)
- Proportional separator lines
- Professional comment windows
- Perfect Turkish text
- 100% content visibility
"""

import json
import math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ==========================================
# CONFIGURATION
# ==========================================
BRICK_RED = RGBColor(192, 57, 43)
DARK_GRAY = RGBColor(44, 62, 80)
LIGHT_BG = RGBColor(245, 247, 250)
BORDER_COLOR = RGBColor(210, 210, 210)
FOOTER_GRAY = RGBColor(150, 150, 150)

# Font sizes
TITLE_SIZE = 28
SUBHEADER_SIZE = 16
BULLET_SIZE = 14
BODY_SIZE = 14
EMPHASIS_SIZE = 14
COMMENT_SIZE = 13

# Layout
MAX_CONTENT_LINES = 14  # Max lines per slide before split
SLIDE_WIDTH = 10.0
CONTENT_LEFT = 0.5
CONTENT_WIDTH = 9.0

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH)
    prs.slide_height = Inches(7.5)
    
    with open('structured_data_final_v7.json', 'r', encoding='utf-8') as f:
        slides_data = json.load(f)
    
    # ==========================================
    # TITLE SLIDE
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Main title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.3), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "OTTOBITE"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = BRICK_RED
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    
    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3.5), Inches(3.35), Inches(3), Pt(4)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BRICK_RED
    line.line.fill.background()
    
    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.6), Inches(8), Inches(1))
    p2 = txBox2.text_frame.add_paragraph()
    p2.text = "Garson Davranışları ve İş Önceliği Rehberi"
    p2.font.size = Pt(24)
    p2.font.color.rgb = DARK_GRAY
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = "Calibri Light"
    
    # ==========================================
    # CONTENT SLIDES
    # ==========================================
    def estimate_lines(item):
        """Estimate how many lines an item will take"""
        text_len = len(item['text'])
        chars_per_line = 75
        lines = max(1, math.ceil(text_len / chars_per_line))
        
        if item['type'] == 'subheader':
            lines += 0.5  # Extra spacing
        elif item['type'] == 'intro_box':
            lines += 0.3
        
        return lines
    
    def add_content_slide(prs, title, content_items, is_continuation=False):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        display_title = f"{title} (Devam)" if is_continuation else title
        
        title_box = slide.shapes.add_textbox(
            Inches(CONTENT_LEFT), Inches(0.35),
            Inches(CONTENT_WIDTH), Inches(0.8)
        )
        tp = title_box.text_frame.add_paragraph()
        tp.text = display_title
        tp.font.size = Pt(TITLE_SIZE)
        tp.font.bold = True
        tp.font.color.rgb = BRICK_RED
        tp.font.name = "Arial"
        
        # Dynamic separator line (proportional to title)
        char_count = len(display_title)
        line_width = min(CONTENT_WIDTH, max(1.5, char_count * 0.18))
        
        separator = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(CONTENT_LEFT), Inches(1.15),
            Inches(line_width), Pt(3)
        )
        separator.fill.solid()
        separator.fill.fore_color.rgb = BRICK_RED
        separator.line.fill.background()
        
        # Separate intro items from main content
        intro_items = [i for i in content_items if i['type'] == 'intro_box']
        main_items = [i for i in content_items if i['type'] != 'intro_box']
        
        current_top = 1.4
        
        # Comment window for intro items
        if intro_items:
            total_intro_chars = sum(len(i['text']) for i in intro_items)
            num_intro = len(intro_items)
            est_lines = max(num_intro, math.ceil(total_intro_chars / 80))
            box_height = min(1.6, max(0.5, est_lines * 0.28 + 0.25))
            
            comment_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(CONTENT_LEFT), Inches(current_top),
                Inches(CONTENT_WIDTH), Inches(box_height)
            )
            comment_box.fill.solid()
            comment_box.fill.fore_color.rgb = LIGHT_BG
            comment_box.line.color.rgb = BORDER_COLOR
            comment_box.line.width = Pt(1.2)
            
            tf = comment_box.text_frame
            tf.margin_left = Inches(0.2)
            tf.margin_right = Inches(0.2)
            tf.margin_top = Inches(0.12)
            tf.margin_bottom = Inches(0.12)
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP
            
            for idx, item in enumerate(intro_items):
                p = tf.add_paragraph()
                p.text = item['text']
                p.font.size = Pt(COMMENT_SIZE)
                p.font.italic = True
                p.font.color.rgb = RGBColor(60, 60, 60)
                p.font.name = "Calibri"
                if idx > 0:
                    p.space_before = Pt(6)
            
            current_top += box_height + 0.2
        
        # Main content
        if main_items:
            available_height = 6.8 - current_top
            
            content_box = slide.shapes.add_textbox(
                Inches(CONTENT_LEFT), Inches(current_top),
                Inches(CONTENT_WIDTH), Inches(available_height)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP
            
            for item in main_items:
                text_val = item['text']
                item_type = item['type']
                
                p = tf.add_paragraph()
                p.font.name = "Calibri"
                
                if item_type == 'subheader':
                    p.text = text_val
                    p.font.bold = True
                    p.font.size = Pt(SUBHEADER_SIZE)
                    p.font.color.rgb = DARK_GRAY
                    p.space_before = Pt(14)
                    p.space_after = Pt(4)
                    
                elif item_type == 'bullet':
                    p.text = "• " + text_val
                    p.font.size = Pt(BULLET_SIZE)
                    p.font.color.rgb = DARK_GRAY
                    p.space_after = Pt(4)
                    
                elif item_type == 'emphasis':
                    p.text = text_val.upper()
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(180, 0, 0)
                    p.font.size = Pt(EMPHASIS_SIZE)
                    p.space_before = Pt(8)
                    p.space_after = Pt(6)
                    
                else:  # body_text
                    p.text = text_val
                    p.font.size = Pt(BODY_SIZE)
                    p.font.color.rgb = DARK_GRAY
                    p.space_after = Pt(6)
        
        # Footer
        footer_box = slide.shapes.add_textbox(
            Inches(8.3), Inches(7.0),
            Inches(1.5), Inches(0.4)
        )
        fp = footer_box.text_frame.add_paragraph()
        fp.text = "OTTOBITE 2026"
        fp.font.size = Pt(9)
        fp.font.color.rgb = FOOTER_GRAY
        fp.alignment = PP_ALIGN.RIGHT
        fp.font.name = "Calibri"
    
    # Process each slide
    for slide_info in slides_data:
        if len(slide_info['content']) == 0:
            continue
        if slide_info['title'] == "OTTOBITE Garson Rehberi" and len(slide_info['content']) < 3:
            continue
        
        title = slide_info['title']
        all_content = slide_info['content']
        
        # Calculate total lines
        total_lines = sum(estimate_lines(item) for item in all_content)
        
        if total_lines <= MAX_CONTENT_LINES:
            # Fits on one slide
            add_content_slide(prs, title, all_content, is_continuation=False)
        else:
            # Need to split
            chunks = []
            current_chunk = []
            current_lines = 0
            
            for item in all_content:
                item_lines = estimate_lines(item)
                
                if current_lines + item_lines > MAX_CONTENT_LINES and current_chunk:
                    chunks.append(current_chunk)
                    current_chunk = []
                    current_lines = 0
                
                current_chunk.append(item)
                current_lines += item_lines
            
            if current_chunk:
                chunks.append(current_chunk)
            
            # Create slides for each chunk
            for idx, chunk in enumerate(chunks):
                add_content_slide(prs, title, chunk, is_continuation=(idx > 0))
    
    # Save
    output_path = "OTTOBITE_SUNUM_PRODUCTION.pptx"
    prs.save(output_path)
    print(f"✓ Production-ready presentation saved: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
