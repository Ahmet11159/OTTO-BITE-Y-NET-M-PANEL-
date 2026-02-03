import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re

def create_bar_presentation_final():
    # --- PREMİUM RENKLER ---
    TERRACOTTA_BG = RGBColor(186, 120, 93)  
    BURGUNDY_BOX = RGBColor(180, 70, 70)    
    DARK_BROWN = RGBColor(50, 35, 30)       
    WHITE = RGBColor(255, 255, 255)
    OFF_WHITE = RGBColor(245, 245, 245)
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def set_bg(slide):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = TERRACOTTA_BG

    content_file = "/Users/ahmet11159/.gemini/antigravity/scratch/bar_sunum_icerik_taslak.md"
    if not os.path.exists(content_file): return

    with open(content_file, 'r', encoding='utf-8') as f:
        full_text = f.read()

    pages_data = re.split(r'---|\n## SAYFA \d+:', full_text)
    
    page_num = 1
    total_slides = 36

    for i, raw_content in enumerate(pages_data):
        content = raw_content.strip()
        if not content: continue
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide)
        
        lines = content.split('\n')
        
        # --- KAPAK SAYFASI (Slide 1) ---
        if i == 0:
            # OTTOBITE
            title_box = slide.shapes.add_textbox(Inches(0), Inches(2.6), prs.slide_width, Inches(1.5))
            tf = title_box.text_frame
            tf.text = "OTTOBITE"
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(90)
            p.font.bold = True
            p.font.color.rgb = DARK_BROWN
            
            # İnce Kırmızı Bar
            bar_w = Inches(7)
            left = (prs.slide_width - bar_w) / 2
            bar = slide.shapes.add_shape(1, left, Inches(4.3), bar_w, Inches(0.4))
            bar.fill.solid()
            bar.fill.fore_color.rgb = BURGUNDY_BOX
            bar.line.width = 0
            
            bar_tf = bar.text_frame
            bar_tf.text = "Bar Departmanı Davranış ve Operasyon Rehberi"
            p2 = bar_tf.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            p2.font.size = Pt(18)
            p2.font.bold = True
            p2.font.color.rgb = WHITE
            
            page_num += 1
            continue

        # --- İÇ SAYFA YERLEŞİM AYARLARI ---
        margin_left = Inches(0.8)
        content_width = prs.slide_width - (margin_left * 2)
        y_pos = Inches(1.5)
        max_y = prs.slide_height - Inches(1.0) # Alttan 1 inç boşluk
        
        # Slayttaki eleman sayısına göre font küçültme
        total_items = len([l for l in lines if l.strip()])
        font_base = 22 if total_items < 10 else 18
        box_font = 20 if total_items < 10 else 17

        for line in lines:
            line = line.strip()
            if not line: continue
            
            # Eğer y_pos sınırı aşıyorsa sayfayı dikeyde sıkıştır (paddingi azalt)
            line_spacing = Inches(0.4) if total_items < 12 else Inches(0.3)

            # Başlık
            if line.startswith('##'):
                title = re.sub(r'^##\s*(\d+\.\s*SAYFA:\s*|SAYFA\s*\d+:\s*)?', '', line)
                tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), content_width, Inches(1))
                tf = tx.text_frame
                tf.text = title
                p = tf.paragraphs[0]
                p.font.size = Pt(36)
                p.font.bold = True
                p.font.color.rgb = DARK_BROWN
                continue

            # Kutulu Vurgu Metni (Kırmızı Kutu)
            if "**" in line:
                box_text = line.replace('**', '')
                # Kutu yüksekliğini metin uzunluğuna göre ayarla
                box_h = Inches(0.5) if len(box_text) < 80 else Inches(0.8)
                
                # Taşma kontrolü
                if y_pos + box_h > max_y: y_pos -= Inches(0.2) # Sıkıştır
                
                box = slide.shapes.add_shape(1, margin_left, y_pos, content_width, box_h)
                box.fill.solid()
                box.fill.fore_color.rgb = BURGUNDY_BOX
                box.line.width = 0
                
                tf = box.text_frame
                tf.word_wrap = True
                tf.text = box_text
                p = tf.paragraphs[0]
                p.font.size = Pt(box_font)
                p.font.color.rgb = WHITE
                y_pos += box_h + Inches(0.3)

            # Liste Elemanları
            elif line.startswith(('•', '*', '-')):
                tx = slide.shapes.add_textbox(margin_left + Inches(0.3), y_pos, content_width - Inches(0.3), Inches(0.4))
                tf = tx.text_frame
                tf.word_wrap = True
                tf.text = line
                p = tf.paragraphs[0]
                p.font.size = Pt(font_base)
                p.font.color.rgb = OFF_WHITE
                y_pos += line_spacing

            # Alt Başlıklar
            elif line.startswith('//'):
                clean = line.replace('//', '').strip()
                tx = slide.shapes.add_textbox(margin_left, y_pos, content_width, Inches(0.5))
                tf = tx.text_frame
                tf.text = clean
                p = tf.paragraphs[0]
                p.font.size = Pt(font_base + 4)
                p.font.bold = True
                p.font.color.rgb = OFF_WHITE
                y_pos += line_spacing + Inches(0.2)

            # Diğer Metinler
            else:
                tx = slide.shapes.add_textbox(margin_left, y_pos, content_width, Inches(0.4))
                tf = tx.text_frame
                tf.word_wrap = True
                tf.text = line
                p = tf.paragraphs[0]
                p.font.size = Pt(font_base)
                p.font.color.rgb = OFF_WHITE
                y_pos += line_spacing

        # Sayfa Numarası
        num_box = slide.shapes.add_textbox(prs.slide_width - Inches(1.5), prs.slide_height - Inches(0.6), Inches(1), Inches(0.3))
        tf = num_box.text_frame
        tf.text = f"{page_num}/{total_slides}"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(140, 110, 90)
        
        page_num += 1

    out = "/Users/ahmet11159/Desktop/OTTOBITE_Bar_Departmani_Rehberi_Final.pptx"
    prs.save(out)
    print(f"Sunum HATASIZ ve PREMİUM olarak hazırlandı: {out}")

if __name__ == "__main__":
    create_bar_presentation_final()
