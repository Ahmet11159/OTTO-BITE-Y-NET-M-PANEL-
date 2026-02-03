# encoding: utf-8
"""
OTTOBITE Final Presentation Generator
Based on detailed PDF analysis - Production Ready
"""

import json
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
BRICK_RED = RGBColor(192, 57, 43)
DARK_GRAY = RGBColor(44, 62, 80)
LIGHT_BG = RGBColor(245, 247, 250)
BORDER_COLOR = RGBColor(215, 215, 215)
FOOTER_GRAY = RGBColor(150, 150, 150)

def create_presentation():
    prs = Presentation()
    
    with open('structured_data_perfect.json', 'r', encoding='utf-8') as f:
        slides_data = json.load(f)
    
    # ==========================================
    # TITLE SLIDE
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    p = txBox.text_frame.add_paragraph()
    p.text = "OTTOBITE"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = BRICK_RED
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1))
    p2 = txBox2.text_frame.add_paragraph()
    p2.text = "Garson Davranışları ve İş Önceliği Rehberi"
    p2.font.size = Pt(28)
    p2.font.color.rgb = DARK_GRAY
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = "Calibri Light"
    
    # Decorative Line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(3.6), Inches(4), Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRICK_RED
    shape.line.fill.background()
    
    # ==========================================
    # CONTENT SLIDES
    # ==========================================
    MAX_LINES = 15
    
    def estimate_lines(items):
        total = 0
        for item in items:
            text_len = len(item['text'])
            lines = max(1, math.ceil(text_len / 70))
            if item['type'] == 'subheader':
                lines += 0.5
            elif item['type'] == 'intro_box':
                lines += 0.3
            total += lines
        return total
    
    def add_slide(prs, title, content, is_continuation=False):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        display_title = f"{title} (Devam)" if is_continuation else title
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
        tp = title_box.text_frame.add_paragraph()
        tp.text = display_title
        tp.font.size = Pt(30)
        tp.font.bold = True
        tp.font.color.rgb = BRICK_RED
        tp.font.name = "Arial"
        
        # Dynamic separator line
        char_count = len(display_title)
        line_width = min(9.0, max(2.0, char_count * 0.20))
        
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.25),
            Inches(line_width), Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = BRICK_RED
        line.line.fill.background()
        
        # Separate intro_box from rest
        intro_items = [i for i in content if i['type'] == 'intro_box']
        main_items = [i for i in content if i['type'] != 'intro_box']
        
        current_top = 1.5
        
        # Comment window for intro items
        if intro_items:
            total_chars = sum(len(i['text']) for i in intro_items)
            num_items = len(intro_items)
            est_lines = max(num_items, math.ceil(total_chars / 80))
            box_height = min(1.5, max(0.5, est_lines * 0.25 + 0.2))
            
            comment_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(current_top),
                Inches(9), Inches(box_height)
            )
            comment_box.fill.solid()
            comment_box.fill.fore_color.rgb = LIGHT_BG
            comment_box.line.color.rgb = BORDER_COLOR
            comment_box.line.width = Pt(1)
            
            tf = comment_box.text_frame
            tf.margin_left = Inches(0.15)
            tf.margin_right = Inches(0.15)
            tf.margin_top = Inches(0.1)
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP
            
            for idx, item in enumerate(intro_items):
                p = tf.add_paragraph()
                p.text = item['text']
                p.font.size = Pt(14)
                p.font.italic = True
                p.font.color.rgb = RGBColor(60, 60, 60)
                p.font.name = "Calibri"
                if idx > 0:
                    p.space_before = Pt(4)
            
            current_top += box_height + 0.15
        
        # Main content
        if main_items:
            available_height = 6.8 - current_top
            
            # Calculate font size
            est_height = estimate_lines(main_items) * 0.35
            if est_height > available_height:
                font_size = max(11, int(16 * (available_height / est_height)))
            else:
                font_size = 16
            
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(current_top),
                Inches(9), Inches(available_height)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP
            
            for item in main_items:
                text_val = item['text']
                item_type = item['type']
                
                p = tf.add_paragraph()
                p.font.name = "Calibri"
                
                if item_type == "subheader":
                    p.text = text_val
                    p.font.bold = True
                    p.font.size = Pt(min(font_size + 2, 18))
                    p.font.color.rgb = DARK_GRAY
                    p.space_before = Pt(10)
                    p.space_after = Pt(3)
                    
                elif item_type == "bullet":
                    p.text = "• " + text_val
                    p.font.size = Pt(font_size)
                    p.font.color.rgb = DARK_GRAY
                    p.space_after = Pt(3)
                    
                elif item_type == "emphasis":
                    p.text = text_val.upper()
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(180, 0, 0)
                    p.font.size = Pt(font_size)
                    p.space_before = Pt(6)
                    p.space_after = Pt(4)
                    
                else:  # body_text
                    p.text = text_val
                    p.font.size = Pt(font_size)
                    p.font.color.rgb = DARK_GRAY
                    p.space_after = Pt(4)
        
        # Footer
        footer_box = slide.shapes.add_textbox(Inches(8), Inches(7.1), Inches(2), Inches(0.5))
        fp = footer_box.text_frame.add_paragraph()
        fp.text = "OTTOBITE 2026"
        fp.font.size = Pt(10)
        fp.font.color.rgb = FOOTER_GRAY
        fp.alignment = PP_ALIGN.RIGHT
    
    # Process slides
    for slide_info in slides_data:
        if len(slide_info['content']) == 0:
            continue
        
        # Skip intro duplicate
        if slide_info['title'] == "OTTOBITE Garson Rehberi":
            continue
        
        title = slide_info['title']
        all_content = slide_info['content']
        
        total_lines = estimate_lines(all_content)
        
        if total_lines <= MAX_LINES:
            add_slide(prs, title, all_content, is_continuation=False)
        else:
            # Split content
            chunks = []
            current_chunk = []
            current_lines = 0
            
            for item in all_content:
                text_len = len(item['text'])
                item_lines = max(1, math.ceil(text_len / 70))
                if item['type'] == 'subheader':
                    item_lines += 0.5
                
                if current_lines + item_lines > MAX_LINES and current_chunk:
                    chunks.append(current_chunk)
                    current_chunk = []
                    current_lines = 0
                
                current_chunk.append(item)
                current_lines += item_lines
            
            if current_chunk:
                chunks.append(current_chunk)
            
            for idx, chunk in enumerate(chunks):
                add_slide(prs, title, chunk, is_continuation=(idx > 0))
    
    # Save
    output = "OTTOBITE_SUNUM_FINAL.pptx"
    prs.save(output)
    print(f"✓ Final presentation saved: {output}")
    return output

if __name__ == "__main__":
    create_presentation()
